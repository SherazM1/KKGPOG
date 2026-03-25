"""
Streamlit Planogram Generator

- Standard Flat Display — all sides on one wide page (original behaviour)
- Full Pallet Display — template-style rebuild (NO raster background), one page per side
"""

from __future__ import annotations

import difflib
import io
import math
import os
import re
from dataclasses import dataclass
from datetime import date
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from openpyxl import load_workbook
import fitz  # PyMuPDF
import numpy as np
import pandas as pd
import pdfplumber
import streamlit as st
from PIL import Image
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas

# ──────────────────────────────────────────────────────────────────────────────
# Constants
# ──────────────────────────────────────────────────────────────────────────────

LAST5_RE = re.compile(r"\b(\d{5})\b")
DIGITS_RE = re.compile(r"\D+")
N_COLS = 3  # standard flat display only

NAVY_RGB = (0.10, 0.16, 0.33)
DISPLAY_STANDARD = "Standard Flat Display"
DISPLAY_FULL_PALLET = "Full Pallet / Multi-Zone Display"

IMAGE_ANCHOR_ROW_0BASED = 5

# ──────────────────────────────────────────────────────────────────────────────
# Data classes
# ──────────────────────────────────────────────────────────────────────────────


@dataclass(frozen=True)
class MatrixRow:
    upc12: str
    norm_name: str
    display_name: str
    cpp_qty: Optional[int]  # from CPP column in Excel — authoritative qty


@dataclass(frozen=True)
class CellData:
    row: int
    col: int
    bbox: Tuple[float, float, float, float]  # (x0, top, x1, bottom) pdfplumber coords
    name: str
    last5: str
    qty: Optional[int]  # label text qty (standard display only; full-pallet uses cpp_qty)
    upc12: Optional[str]


@dataclass(frozen=True)
class PageData:
    # standard display
    page_index: int
    x_bounds: np.ndarray
    y_bounds: np.ndarray
    cells: List[CellData]


@dataclass(frozen=True)
class AnnotationBox:
    kind: str  # bonus_strip | gift_card_holders | marketing_signage | fraud_signage | wm_new_pkg
    label: str
    bbox: Tuple[float, float, float, float]  # pdfplumber coords


@dataclass(frozen=True)
class FullPalletPage:
    page_index: int
    side_letter: str
    cells: List[CellData]
    annotations: List[AnnotationBox]


@dataclass(frozen=True)
class PptCard:
    card_id: str
    title: str
    image_bytes: Optional[bytes] = None
    image_ext: Optional[str] = None


@dataclass(frozen=True)
class PptSideCards:
    side: str
    top8: List[PptCard]
    side6: List[PptCard]


@dataclass(frozen=True)
class GiftHolder:
    side: str
    item_no: str
    name: str
    qty: Optional[int]
    image_bytes: Optional[bytes] = None
    image_ext: Optional[str] = None


# ──────────────────────────────────────────────────────────────────────────────
# Fonts
# ──────────────────────────────────────────────────────────────────────────────


def _safe_register_font(name: str, rel_path: str) -> None:
    try:
        full = str(Path(__file__).resolve().parent / "assets" / rel_path)
        if os.path.isfile(full):
            pdfmetrics.registerFont(TTFont(name, full))
    except Exception:
        pass


_safe_register_font("Raleway", "Raleway-Regular.ttf")
_safe_register_font("Raleway-Bold", "Raleway-Bold.ttf")


def _font(preferred: str, fallback: str) -> str:
    try:
        pdfmetrics.getFont(preferred)
        return preferred
    except Exception:
        return fallback


TITLE_FONT = _font("Raleway-Bold", "Helvetica-Bold")
BODY_FONT = _font("Raleway", "Helvetica")
BODY_BOLD_FONT = _font("Raleway-Bold", "Helvetica-Bold")

# ──────────────────────────────────────────────────────────────────────────────
# Small utilities
# ──────────────────────────────────────────────────────────────────────────────


def _norm_name(s: str) -> str:
    s = (s or "").upper()
    s = re.sub(r"[^A-Z0-9]+", " ", s)
    return re.sub(r"\s+", " ", s).strip()


def _coerce_upc12(v: object) -> Optional[str]:
    if v is None or (isinstance(v, float) and math.isnan(v)):
        return None
    s = re.sub(r"\.0$", "", str(v).strip())
    s = DIGITS_RE.sub("", s)
    return s.zfill(12) if s else None


def _coerce_int(v: object) -> Optional[int]:
    if v is None or (isinstance(v, float) and math.isnan(v)):
        return None
    s = re.sub(r"[^\d\-]", "", re.sub(r"\.0$", "", str(v).strip()))
    try:
        return int(s) if s else None
    except ValueError:
        return None


def _to_last5(v: object) -> str:
    s = DIGITS_RE.sub("", str(v or "").strip())
    if not s:
        return ""
    if len(s) >= 5:
        return s[-5:]
    return s.zfill(5)


def _norm_header(v: object) -> str:
    s = re.sub(r"[^A-Z0-9]+", "_", str(v).strip().upper()).strip("_")
    return s or "COL"


def _find_header_row(df: pd.DataFrame) -> int:
    for i in range(min(len(df), 50)):
        row = df.iloc[i].astype(str).str.upper().tolist()
        if any("UPC" in c for c in row) and any(
            tok in " ".join(row) for tok in ("NAME", "DESCRIPTION", "CPP")
        ):
            return i
    return 0


def _pick_col(cols: List[str], tokens: List[str], fallback: int) -> str:
    uc = [c.upper() for c in cols]
    for t in tokens:
        for i, c in enumerate(uc):
            if t in c:
                return cols[i]
    return cols[min(fallback, len(cols) - 1)]


def _pick_col_optional(cols: List[str], tokens: List[str]) -> Optional[str]:
    uc = [c.upper() for c in cols]
    for t in tokens:
        for i, c in enumerate(uc):
            if t in c:
                return cols[i]
    return None


def _hex_to_rgb(h: str) -> Tuple[float, float, float]:
    h = (h or "").lstrip("#").strip()
    if len(h) != 6:
        return NAVY_RGB
    return int(h[:2], 16) / 255, int(h[2:4], 16) / 255, int(h[4:], 16) / 255


def _try_load_logo() -> Optional[Image.Image]:
    for p in [
        Path.cwd() / "assets" / "KKG-Logo-02.png",
        Path(__file__).resolve().parent / "assets" / "KKG-Logo-02.png",
    ]:
        if p.exists():
            try:
                return Image.open(p).convert("RGBA")
            except Exception:
                pass
    return None


# ──────────────────────────────────────────────────────────────────────────────
# Matrix loading
# ──────────────────────────────────────────────────────────────────────────────


@st.cache_data(show_spinner=False)
def load_matrix_index(matrix_bytes: bytes) -> Dict[str, List[MatrixRow]]:
    df_raw = pd.read_excel(io.BytesIO(matrix_bytes), header=None)

    hrow = _find_header_row(df_raw)

    headers: List[str] = []
    seen: Dict[str, int] = {}
    for v in df_raw.iloc[hrow].tolist():
        base = _norm_header(v)
        n = seen.get(base, 0)
        seen[base] = n + 1
        headers.append(base if n == 0 else f"{base}_{n+1}")

    df = df_raw.iloc[hrow + 1 :].copy()
    df.columns = headers

    upc_col = _pick_col(headers, ["UPC"], 0)
    name_col = _pick_col(headers, ["NAME", "DESCRIPTION"], 1 if len(headers) > 1 else 0)
    cpp_col = _pick_col_optional(headers, ["CPP"])

    df["__upc12"] = df[upc_col].map(_coerce_upc12)
    df["__name"] = df[name_col].astype(str).fillna("")

    if cpp_col and cpp_col in df.columns:
        df["__cpp"] = df[cpp_col].map(_coerce_int)
    else:
        df["__cpp"] = None

    df = df[df["__upc12"].notna()].copy()
    df["__last5"] = df["__upc12"].str[-5:]
    df["__norm"] = df["__name"].map(_norm_name)

    idx: Dict[str, List[MatrixRow]] = {}
    for _, r in df.iterrows():
        last5 = str(r["__last5"])
        cpp_val = r.get("__cpp")
        cpp = None if pd.isna(cpp_val) else int(cpp_val) if cpp_val is not None else None
        display_name = str(r["__name"]).strip()
        idx.setdefault(last5, []).append(
            MatrixRow(
                upc12=str(r["__upc12"]),
                norm_name=str(r["__norm"]),
                display_name=display_name,
                cpp_qty=cpp,
            )
        )
    return idx


@st.cache_data(show_spinner=False)
def load_full_pallet_matrix_index(matrix_bytes: bytes) -> Dict[str, List[MatrixRow]]:
    """Full-pallet matrix index.

    Index key  = last 5 digits of the 11-digit UPC column (column header "UPC").
    Stored upc12 = 11-digit UPC zero-padded to 12 chars.

    The labels PDF encodes each product as the last 5 digits of its 11-digit UPC,
    so indexing on that column is required for correct resolution.  The zero-padded
    11-digit value matches what the reference planogram displays as the product UPC.
    """
    df_raw = pd.read_excel(io.BytesIO(matrix_bytes), header=None)
    hrow = _find_header_row(df_raw)

    headers: List[str] = []
    seen: Dict[str, int] = {}
    for v in df_raw.iloc[hrow].tolist():
        base = _norm_header(v)
        n = seen.get(base, 0)
        seen[base] = n + 1
        headers.append(base if n == 0 else f"{base}_{n+1}")

    df = df_raw.iloc[hrow + 1 :].copy()
    df.columns = headers

    name_col = _pick_col(headers, ["NAME", "DESCRIPTION"], 1 if len(headers) > 1 else 0)
    cpp_col = _pick_col_optional(headers, ["CPP"])

    # ── FIX A: use ONLY the 11-digit UPC column as the index key ─────────────
    # The labels PDF encodes items as the last 5 digits of the 11-digit UPC.
    # Using the 12-digit UPC column would give a different last-5 (check digit
    # appended) and break all lookups.  Prefer the exact header "UPC"; fall back
    # to the first column whose normalised name contains "UPC".
    upc11_col: Optional[str] = None
    for h in headers:
        if h.upper() == "UPC":
            upc11_col = h
            break
    if upc11_col is None:
        upc11_col = _pick_col(headers, ["UPC"], 0)

    df["__name"] = df[name_col].astype(str).fillna("")
    if cpp_col and cpp_col in df.columns:
        df["__cpp"] = df[cpp_col].map(_coerce_int)
    else:
        df["__cpp"] = None
    df["__norm"] = df["__name"].map(_norm_name)

    idx: Dict[str, List[MatrixRow]] = {}
    seen_pairs: set = set()
    for _, r in df.iterrows():
        display_name = str(r["__name"]).strip()
        cpp_val = r.get("__cpp")
        cpp = None if pd.isna(cpp_val) else int(cpp_val) if cpp_val is not None else None

        raw_upc11 = r.get(upc11_col)
        if raw_upc11 is None:
            continue
        upc11_str = re.sub(r"\.0$", "", str(raw_upc11).strip())
        digits11 = re.sub(r"[^0-9]", "", upc11_str)
        if len(digits11) < 5:
            continue
        last5 = digits11[-5:]
        # Store as 11-digit UPC zero-padded to 12 chars — this is the format
        # shown in the reference planogram output (e.g. "019674209969").
        upc12_display = digits11.zfill(12)

        pair = (last5, upc12_display)
        if pair in seen_pairs:
            continue
        seen_pairs.add(pair)
        idx.setdefault(last5, []).append(
            MatrixRow(
                upc12=upc12_display,
                norm_name=str(r["__norm"]),
                display_name=display_name,
                cpp_qty=cpp,
            )
        )
    return idx


def resolve_full_pallet(last5: str, label_name: str, idx: Dict[str, List[MatrixRow]]) -> Optional[MatrixRow]:
    key = _to_last5(last5)
    rows = idx.get(key, [])
    if not rows:
        return None
    if len(rows) == 1:
        return rows[0]
    target = _norm_name(label_name)
    return max(rows, key=lambda r: difflib.SequenceMatcher(None, target, r.norm_name).ratio())


def _resolve(last5: str, label_name: str, idx: Dict[str, List[MatrixRow]]) -> Optional[MatrixRow]:
    rows = idx.get(last5, [])
    if not rows:
        return None
    if len(rows) == 1:
        return rows[0]
    target = _norm_name(label_name)
    return max(
        rows, key=lambda r: difflib.SequenceMatcher(None, target, r.norm_name).ratio()
    )


def _coerce_item_no(v: object) -> Optional[str]:
    if v is None or (isinstance(v, float) and math.isnan(v)):
        return None
    s = re.sub(r"\.0$", "", str(v).strip())
    s = DIGITS_RE.sub("", s)
    return s.zfill(6) if s else None


@st.cache_data(show_spinner=False)
def load_ppt_cards(pptx_bytes: bytes) -> Dict[str, PptSideCards]:
    """Parse PPTX per-side slides and return PptSideCards with images.

    Structure of each SIDE slide (slides 2-5 = A-D):
      - Top-level GROUP shapes (one per card image, width < 10in) sitting in the top ~1in band
      - TEXT_BOX shapes with "ID #nn" labels positioned just below each group
      - One loose PICTURE shape may appear on slides where one card isn't inside a group

    Each card GROUP contains:
      - The actual card image (variable blob size)
      - The Walmart watermark/logo (constant 71369-byte blob) — skip this

    Algorithm:
      1. Collect all image containers (groups + loose pictures), extract card blob
      2. Collect all TEXT_BOX labels with ID#
      3. Match each label to the nearest image container above it (by cx proximity)
      4. Split into top8 (labels at y < 40% slide height) and side6 (y >= 40%)
      5. Sort top8 left→right by cx; sort side6 by (row-band, cx)
    """
    from pptx import Presentation  # type: ignore
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore

    WATERMARK_BLOB_SIZE = 71369  # Walmart logo — same bytes across all groups

    prs = Presentation(io.BytesIO(pptx_bytes))
    slide_h = float(prs.slide_height)
    slide_w = float(prs.slide_width)

    id_re = re.compile(r"\bID\s*#?\s*[:\-]?\s*(\d{1,8})\b", re.IGNORECASE)
    side_re = re.compile(r"\bSIDE\s*([A-D])\b", re.IGNORECASE)

    def _extract_card_pictures(sh) -> List[dict]:
        """Return the largest non-watermark image blob from a shape or group."""
        pics: List[dict] = []

        def _collect(shapes) -> None:
            for s in shapes:
                st = getattr(s, "shape_type", None)
                if st == MSO_SHAPE_TYPE.PICTURE:
                    b = bytes(s.image.blob)
                    if len(b) == WATERMARK_BLOB_SIZE:
                        continue
                    l = float(getattr(s, "left", 0) or 0)
                    t = float(getattr(s, "top", 0) or 0)
                    w = float(getattr(s, "width", 0) or 0)
                    h = float(getattr(s, "height", 0) or 0)
                    pics.append({
                        "blob": b,
                        "cx": l + w / 2,
                        "cy": t + h / 2,
                        "area": w * h,
                        "byte_size": len(b)
                    })

                elif st == MSO_SHAPE_TYPE.GROUP:
                    _collect(s.shapes)

        st = getattr(sh, "shape_type", None)
        if st == MSO_SHAPE_TYPE.PICTURE:
            b = bytes(sh.image.blob)
            if len(b) != WATERMARK_BLOB_SIZE:
                l = float(getattr(sh, "left", 0) or 0)
                t = float(getattr(sh, "top", 0) or 0)
                w = float(getattr(sh, "width", 0) or 0)
                h = float(getattr(sh, "height", 0) or 0)
                pics.append({
                    "blob": b,
                    "cx": l + w / 2,
                    "cy": t + h / 2,
                    "area": w * h,
                    "byte_size": len(b),
                })
        elif st == MSO_SHAPE_TYPE.GROUP:
            _collect(sh.shapes)

        return pics

    best_by_side: Dict[str, Tuple[int, List[PptCard], List[PptCard]]] = {}

    for slide in prs.slides:
        # ── Detect side letter ──────────────────────────────────────────────
        side_letter: Optional[str] = None
        for sh in slide.shapes:
            m = side_re.search(str(getattr(sh, "text", "") or ""))
            if m:
                side_letter = m.group(1).upper()
                break
        if side_letter is None:
            continue  # skip title / intro / half-pallet slides

        # ── Collect image containers ────────────────────────────────────────
        img_containers: List[dict] = []
        for sh in slide.shapes:
            stype = getattr(sh, "shape_type", None)
            l = float(getattr(sh, "left", 0) or 0)
            t = float(getattr(sh, "top", 0) or 0)
            w = float(getattr(sh, "width", 0) or 0)
            h = float(getattr(sh, "height", 0) or 0)

            if stype == MSO_SHAPE_TYPE.GROUP:
                if w > slide_w * 0.80:
                    continue  # skip full-slide background frame
                candidates = _extract_card_pictures(sh)
                if not candidates:
                    continue
            elif stype == MSO_SHAPE_TYPE.PICTURE:
                candidates = _extract_card_pictures(sh)
                if not candidates:
                    continue
            else:
                continue

            img_containers.append({
                "cx": l + w / 2,
                "cy": t + h / 2,
                "bottom": t + h,
                "candidates": candidates,
            })

        # ── Collect labels ──────────────────────────────────────────────────
        labels: List[dict] = []
        for sh in slide.shapes:
            txt = str(getattr(sh, "text", "") or "").strip()
            m = id_re.search(txt)
            if not m:
                continue
            card_id = m.group(1)
            lines = [ln.strip() for ln in txt.splitlines() if ln.strip()]
            id_idx = next((i for i, ln in enumerate(lines) if id_re.search(ln)), 0)
            title = " ".join(lines[:id_idx]).strip()
            l = float(getattr(sh, "left", 0) or 0)
            t = float(getattr(sh, "top", 0) or 0)
            w = float(getattr(sh, "width", 0) or 0)
            h = float(getattr(sh, "height", 0) or 0)
            labels.append({
                "card_id": card_id,
                "title": title,
                "cx": l + w / 2,
                "top": t,
                "cy": t + h / 2,
            })

        if not labels:
            continue

        # ── Match each label to the nearest image above it ──────────────────
        # Process in reading order (top→bottom, left→right) so greedy works well.
        used: set = set()
        label_to_img: Dict[str, Optional[dict]] = {}

        for lab in sorted(labels, key=lambda d: (round(d["top"] / (slide_h * 0.05)), d["cx"])):
            best_i: Optional[int] = None
            best_score = float("inf")
            for i, img in enumerate(img_containers):
                if i in used:
                    continue
                # Image must be above (or at most 0.5in below) the label top
                if img["bottom"] > lab["top"] + slide_h * 0.07:
                    continue
                dx = abs(lab["cx"] - img["cx"])
                dy = max(0.0, lab["top"] - img["bottom"])
                score = dx + 0.3 * dy
                if score < best_score:
                    best_score = score
                    best_i = i
            if best_i is not None:
                used.add(best_i)
                label_to_img[lab["card_id"]] = img_containers[best_i]
            else:
                label_to_img[lab["card_id"]] = None

        # ── Split into top8 and side6 by y-position ─────────────────────────
        # Labels for the main top row appear at ~30-36% of slide height.
        # Labels for the side panel appear at ~57-93% of slide height.
        TOP_THRESH = slide_h * 0.42
        top_labels = sorted(
            [lb for lb in labels if lb["top"] < TOP_THRESH], key=lambda d: d["cx"]
        )
        side_labels = sorted(
            [lb for lb in labels if lb["top"] >= TOP_THRESH],
            key=lambda d: (round(d["top"] / (slide_h * 0.12)), d["cx"]),
        )

        def _make_card(lab: dict) -> PptCard:
            img_entry = label_to_img.get(lab["card_id"])
            img_bytes: Optional[bytes] = None
            img_ext: Optional[str] = None
            if img_entry:
                candidates = img_entry.get("candidates", [])
                chosen: Optional[dict] = None

                if candidates:
                    chosen = min(
                        candidates,
                        key=lambda c: (abs(c["cx"] - lab["cx"]), abs(c["cy"] - lab["cy"]), -c["area"], -c["byte_size"]),
                    )
                # Normalise to PNG
                if chosen:
                    raw = chosen["blob"]
                    try:
                        from io import BytesIO as _BIO
                        im = Image.open(_BIO(raw)).convert("RGBA")
                        out = _BIO()
                        im.save(out, format="PNG")
                        img_bytes = out.getvalue()
                        img_ext = "png"
                    except Exception:
                        img_bytes = raw
                        img_ext = "png"

            return PptCard(
                card_id=lab["card_id"],
                title=lab["title"],
                image_bytes=img_bytes,
                image_ext=img_ext,
            )

        top8 = [_make_card(lb) for lb in top_labels[:8]]
        side6 = [_make_card(lb) for lb in side_labels[:6]]

        total = len(top8) + len(side6)
        prev = best_by_side.get(side_letter)
        if prev is None or total > prev[0]:
            best_by_side[side_letter] = (total, top8, side6)

    parsed: Dict[str, PptSideCards] = {}
    for side in "ABCD":
        entry = best_by_side.get(side)
        if entry:
            _, top8, side6 = entry
        else:
            top8, side6 = [], []
        parsed[side] = PptSideCards(side=side, top8=top8, side6=side6)

    return parsed


def validate_ppt_side_cards(ppt_cards: Dict[str, PptSideCards]) -> List[str]:
    issues: List[str] = []
    for side in "ABCD":
        side_cards = ppt_cards.get(side, PptSideCards(side=side, top8=[], side6=[]))
        if len(side_cards.top8) != 8 or len(side_cards.side6) != 6:
            issues.append(
                f"SIDE {side}: found top8={len(side_cards.top8)} and side6={len(side_cards.side6)} (expected 8 and 6)."
            )
    return issues


def _img_anchor_col(img) -> Optional[int]:
    try:
        anchor = getattr(img, "anchor", None)
        if anchor is None:
            return None
        _from = getattr(anchor, "_from", None)
        if _from is not None:
            return int(_from.col)
        frm = getattr(anchor, "from_", None)
        if frm is not None:
            return int(frm.col)
    except Exception:
        return None
    return None


def _img_bytes_and_ext(img) -> Tuple[Optional[bytes], Optional[str]]:
    try:
        raw = None
        if hasattr(img, "_data"):
            raw = img._data()
        elif hasattr(img, "ref"):
            ref = img.ref
            if isinstance(ref, bytes):
                raw = ref
            elif hasattr(ref, "read"):
                raw = ref.read()
        if not raw:
            return None, None

        fmt = str(getattr(img, "format", "") or "").strip().lower()
        ext = fmt if fmt in {"png", "jpg", "jpeg", "gif", "bmp"} else "png"

        try:
            im = Image.open(io.BytesIO(raw)).convert("RGBA")
            out = io.BytesIO()
            im.save(out, format="PNG")
            return out.getvalue(), "png"
        except Exception:
            return raw, ext
    except Exception:
        return None, None


def _extract_ws_images_by_col(ws) -> Dict[int, Tuple[bytes, str]]:
    """
    Map 0-based Excel column index -> first embedded image found in that column.
    """
    out: Dict[int, Tuple[bytes, str]] = {}
    for img in getattr(ws, "_images", []) or []:
        col = _img_anchor_col(img)
        if col is None:
            continue
        img_bytes, img_ext = _img_bytes_and_ext(img)
        if not img_bytes:
            continue
        out.setdefault(col, (img_bytes, img_ext or "png"))
    return out


def _nearest_image_for_col(
    images_by_col: Dict[int, Tuple[bytes, str]],
    target_col: int,
    max_distance: int = 3,
) -> Tuple[Optional[bytes], Optional[str]]:
    best: Optional[Tuple[int, Tuple[bytes, str]]] = None
    for col, payload in images_by_col.items():
        dist = abs(int(col) - int(target_col))
        if dist > max_distance:
            continue
        if best is None or dist < best[0]:
            best = (dist, payload)
    return best[1] if best else (None, None)


def _image_for_col_span(
    images_by_col: Dict[int, Tuple[bytes, str]],
    start_col: int,
    end_col: int,
) -> Tuple[Optional[bytes], Optional[str]]:
    in_span = [
        (col, payload)
        for col, payload in images_by_col.items()
        if start_col <= int(col) <= end_col
    ]
    if in_span:
        # Prefer the left-most image in the slot span; there should normally be one.
        in_span.sort(key=lambda x: x[0])
        return in_span[0][1]
    center_col = (start_col + end_col) // 2
    return _nearest_image_for_col(images_by_col, center_col, max_distance=2)


def _first_nonempty_in_span(ws, row_idx: int, start_col: int, end_col: int) -> Optional[str]:
    for c in range(start_col, end_col + 1):
        val = ws.cell(row=row_idx + 1, column=c + 1).value
        if val is None:
            continue
        text = str(val).strip()
        if text:
            return text
    return None


def _first_int_in_span(ws, row_idx: int, start_col: int, end_col: int) -> Optional[int]:
    for c in range(start_col, end_col + 1):
        val = ws.cell(row=row_idx + 1, column=c + 1).value
        if val is None:
            continue
        text = str(val).strip()
        if not text:
            continue
        digits = re.sub(r"[^\d]", "", text)
        if digits:
            try:
                return int(digits)
            except Exception:
                pass
    return None


def _extract_top_holder_slots(
    ws,
    images_by_col: Dict[int, Tuple[bytes, str]],
    qty_row: int,
    item_row: int,
    desc_row: int,
) -> List[dict]:
    """
    Extract holder slots from the top FULL PALLET block:
    POCKET/PEG headers -> image -> qty/item/description below.
    """
    header_row = IMAGE_ANCHOR_ROW_0BASED
    if header_row < 1:
        return []

    max_col = ws.max_column
    slot_starts: List[Tuple[int, str]] = []

    for cidx in range(max_col):
        raw = ws.cell(row=header_row + 1, column=cidx + 1).value
        text = str(raw).strip() if raw is not None else ""
        text_u = text.upper()
        if text_u.startswith("POCKET ") or text_u.startswith("PEG "):
            slot_starts.append((cidx, text))
        elif "MARKETING MESSAGE PANEL" in text_u:
            slot_starts.append((cidx, "__BREAK__"))

    if not slot_starts:
        return []

    slots: List[dict] = []
    real_headers = [(c, h) for c, h in slot_starts if h != "__BREAK__"]

    for i, (start_col, header_text) in enumerate(real_headers):
        next_boundaries = [
            c for c, _ in slot_starts
            if c > start_col
        ]
        end_col = (min(next_boundaries) - 1) if next_boundaries else (max_col - 1)
        if end_col < start_col:
            end_col = start_col

        item_no = _first_nonempty_in_span(ws, item_row, start_col, end_col)
        if not item_no:
            continue

        qty = _first_int_in_span(ws, qty_row, start_col, end_col)
        name = _first_nonempty_in_span(ws, desc_row, start_col, end_col) or ""

        img_bytes, img_ext = _image_for_col_span(images_by_col, start_col, end_col)

        slots.append(
            {
                "header": header_text,
                "start_col": start_col,
                "end_col": end_col,
                "item_no": re.sub(r"[^\d]", "", str(item_no)),
                "qty": qty,
                "name": name,
                "image_bytes": img_bytes,
                "image_ext": img_ext,
            }
        )

    return slots


@st.cache_data(show_spinner=False)
def load_gift_card_holders(gift_bytes: bytes) -> Dict[str, List[GiftHolder]]:
    """Parse POG workbook holder section and return ordered per-side holders."""

    try:
        xls = pd.read_excel(io.BytesIO(gift_bytes), sheet_name=None, header=None)
    except Exception as e:
        raise ValueError(f"Unable to read Gift Card Holders workbook: {e}")

    def find_header_row(df: pd.DataFrame, tokens: List[str], max_rows: int = 120) -> int:
        for i in range(min(len(df), max_rows)):
            row = df.iloc[i].astype(str).fillna("").str.upper().tolist()
            if all(any(tok in c for c in row) for tok in tokens):
                return i
        return -1

    def normalize_headers(raw_headers: List[object]) -> List[str]:
        headers: List[str] = []
        seen: Dict[str, int] = {}
        for v in raw_headers:
            base = _norm_header(v)
            n = seen.get(base, 0) + 1
            seen[base] = n
            headers.append(base if n == 1 else f"{base}_{n}")
        return headers

    def pick_col(headers: List[str], tokens: List[str]) -> Optional[str]:
        up = [h.upper() for h in headers]
        for tok in tokens:
            for i, h in enumerate(up):
                if tok in h:
                    return headers[i]
        return None

    def parse_side(value: object, current_side: str) -> str:
        text = str(value or "").upper()
        m = re.search(r"SIDE\s*([A-D])", text)
        if m:
            return m.group(1)
        m = re.fullmatch(r"\s*([A-D])\s*", text)
        if m:
            return m.group(1)
        return current_side

    full_pallet_sheet = next(
        (name for name in xls.keys() if "FULL" in name.upper() and "PALLET" in name.upper()),
        None,
    )
    if full_pallet_sheet is None:
        raise ValueError("FULL PALLET sheet/table missing in workbook.")

    try:
        wb = load_workbook(io.BytesIO(gift_bytes), data_only=True)
    except Exception as e:
        raise ValueError(f"Unable to open Gift Card Holders workbook images: {e}")

    full_pallet_ws = None
    for ws in wb.worksheets:
        title_u = str(ws.title or "").upper()
        if "FULL" in title_u and "PALLET" in title_u:
            full_pallet_ws = ws
            break

    if full_pallet_ws is None:
        raise ValueError("Could not find FULL PALLET sheet in workbook.")

    images_by_col = _extract_ws_images_by_col(full_pallet_ws)

    desc_map: Dict[str, str] = {}
    for _, sheet_df in xls.items():
        hdr = find_header_row(sheet_df, ["ITEM", "DESCRIPTION"])
        if hdr < 0:
            continue
        cols = normalize_headers(sheet_df.iloc[hdr].tolist())
        data = sheet_df.iloc[hdr + 1 :].copy()
        data.columns = cols
        item_lookup_col = pick_col(cols, ["ITEM", "ITEM_#", "SKU"])
        desc_lookup_col = pick_col(cols, ["DESCRIPTION", "DESC", "NAME"])
        if not item_lookup_col or not desc_lookup_col:
            continue
        for _, row in data.iterrows():
            item_no = _coerce_item_no(row.get(item_lookup_col))
            if not item_no:
                continue
            desc = str(row.get(desc_lookup_col, "") or "").strip()
            if desc:
                desc_map[item_no] = desc

    full_df_raw = xls[full_pallet_sheet].copy()

    item_row = -1
    qty_row = -1
    desc_row = -1
    for i in range(len(full_df_raw)):
        row_vals = [str(v or "").strip().upper() for v in full_df_raw.iloc[i].tolist()]
        if item_row < 0 and any("ITEM #" in v or v == "ITEM" for v in row_vals):
            numeric_items = sum(1 for v in row_vals if _coerce_item_no(v))
            if numeric_items >= 4:
                item_row = i
        if item_row >= 0 and i <= item_row and qty_row < 0:
            if any(v == "QTY" or v.startswith("QTY") for v in row_vals):
                qty_row = i
        if item_row >= 0 and i >= item_row and desc_row < 0:
            if any("DESCRIPTION" in v for v in row_vals):
                desc_row = i
        if item_row >= 0 and qty_row >= 0 and desc_row >= 0:
            break

    top_slots = _extract_top_holder_slots(
        full_pallet_ws,
        images_by_col=images_by_col,
        qty_row=qty_row,
        item_row=item_row,
        desc_row=desc_row,
    )

    holders: Dict[str, List[GiftHolder]] = {s: [] for s in "ABCD"}
    if item_row >= 0 and qty_row >= 0:
        side_markers: List[Tuple[int, str]] = []
        for r in range(max(0, item_row - 5), item_row + 1):
            vals = [str(v or "").strip().upper() for v in full_df_raw.iloc[r].tolist()]
            for cidx, txt in enumerate(vals):
                m = re.search(r"SIDE\s*([A-D])", txt)
                if m:
                    side_markers.append((cidx, m.group(1)))
        side_markers = sorted(side_markers, key=lambda t: t[0])

        def side_for_col(cidx: int) -> str:
            if not side_markers:
                return "A"
            chosen = side_markers[0][1]
            for sc, ss in side_markers:
                if cidx >= sc:
                    chosen = ss
                else:
                    break
            return chosen

        if top_slots:
            for slot in top_slots:
                cidx = int(slot["start_col"])
                side = side_for_col(cidx)
                item_no = str(slot["item_no"]).strip()
                if not item_no:
                    continue
                name = str(slot["name"] or "").strip()
                qty = slot["qty"]
                img_bytes = slot["image_bytes"]
                img_ext = slot["image_ext"]
                holders.setdefault(side, []).append(
                    GiftHolder(
                        side=side,
                        item_no=item_no,
                        name=name,
                        qty=qty,
                        image_bytes=img_bytes,
                        image_ext=img_ext,
                    )
                )

    if any(holders.values()):
        non_empty = [s for s in "ABCD" if holders.get(s)]
        if non_empty == ["A"]:
            base_list = holders["A"]
            for s in "BCD":
                holders[s] = [
                    GiftHolder(
                        side=s,
                        item_no=h.item_no,
                        name=h.name,
                        qty=h.qty,
                        image_bytes=h.image_bytes,
                        image_ext=h.image_ext,
                    )
                    for h in base_list
                ]
        return holders

    header_row = find_header_row(full_df_raw, ["ITEM", "QTY"])
    if header_row < 0:
        raise ValueError("FULL PALLET holder table missing ITEM/QTY columns.")

    headers = normalize_headers(full_df_raw.iloc[header_row].tolist())
    full_df = full_df_raw.iloc[header_row + 1 :].copy()
    full_df.columns = headers

    item_col = pick_col(headers, ["ITEM", "ITEM_#", "SKU"])
    qty_col = pick_col(headers, ["QTY", "QUANTITY"])
    side_col = pick_col(headers, ["SIDE"])
    if item_col is None or qty_col is None:
        raise ValueError("FULL PALLET holder table missing ITEM/QTY columns.")

    current_side = "A"
    for _, row in full_df.iterrows():
        side_source = row.get(side_col) if side_col else " "
        current_side = parse_side(side_source, current_side)

        item_no = _coerce_item_no(row.get(item_col))
        if not item_no:
            row_text = " ".join(str(v or "") for v in row.tolist())
            current_side = parse_side(row_text, current_side)
            continue

        qty = _coerce_int(row.get(qty_col))
        name = desc_map.get(item_no) or "(missing description)"
        side = current_side if current_side in "ABCD" else "A"
        holders.setdefault(side, []).append(
            GiftHolder(
                side=side,
                item_no=item_no,
                name=name,
                qty=qty,
                image_bytes=None,
                image_ext=None,
            )
        )

    if not any(holders.values()):
        raise ValueError("No holder Item # rows found in FULL PALLET table.")

    non_empty = [s for s in "ABCD" if holders.get(s)]
    if non_empty == ["A"]:
        base_list = holders["A"]
        for s in "BCD":
            holders[s] = [
                GiftHolder(
                    side=s,
                    item_no=h.item_no,
                    name=h.name,
                    qty=h.qty,
                    image_bytes=h.image_bytes,
                    image_ext=h.image_ext,
                )
                for h in base_list
            ]
    return holders


# ──────────────────────────────────────────────────────────────────────────────
# Clustering helpers
# ──────────────────────────────────────────────────────────────────────────────


def kmeans_1d(values: List[float], k: int, iters: int = 40) -> np.ndarray:
    v = np.asarray(values, dtype=float)
    if len(v) == 0:
        return np.array([], dtype=float)
    k = max(1, min(k, len(v)))
    qs = np.linspace(0, 1, k, endpoint=False) + 0.5 / k
    centers = np.quantile(v, qs)

    for _ in range(iters):
        d = np.abs(v[:, None] - centers[None, :])
        labels = d.argmin(axis=1)
        nc = centers.copy()
        for i in range(k):
            pts = v[labels == i]
            if len(pts):
                nc[i] = float(pts.mean())
        if np.allclose(nc, centers):
            break
        centers = nc

    return np.sort(centers)


def cluster_positions(values: List[float], tol: float) -> np.ndarray:
    if not values:
        return np.array([], dtype=float)
    vals = sorted(float(v) for v in values)
    groups = [[vals[0]]]
    for v in vals[1:]:
        if abs(v - float(np.mean(groups[-1]))) <= tol:
            groups[-1].append(v)
        else:
            groups.append([v])
    return np.array([float(np.mean(g)) for g in groups])


def boundaries_from_centers(centers: np.ndarray) -> np.ndarray:
    c = np.sort(np.asarray(centers, dtype=float))
    if len(c) == 0:
        return c
    if len(c) == 1:
        return np.array([c[0] - 50, c[0] + 50])
    mids = (c[:-1] + c[1:]) / 2
    left = c[0] - (mids[0] - c[0])
    right = c[-1] + (c[-1] - mids[-1])
    return np.concatenate([[left], mids, [right]])


# ──────────────────────────────────────────────────────────────────────────────
# Label parsing
# ──────────────────────────────────────────────────────────────────────────────


def parse_label_cell_text(text: str) -> Tuple[str, str, Optional[int]]:
    lines = [ln.strip() for ln in (text or "").splitlines() if ln.strip()]
    joined = " ".join(lines)

    m = LAST5_RE.search(joined)
    last5 = m.group(1) if m else ""

    nums = re.findall(r"\b(\d{1,3})\b", joined)
    qty = int(nums[-1]) if nums else None

    name = " ".join(
        ln
        for ln in lines
        if not (last5 and last5 in ln)
        and not (qty is not None and re.fullmatch(str(qty), ln))
    ).strip()

    return name, last5, qty


# ──────────────────────────────────────────────────────────────────────────────
# Standard display extraction (UNCHANGED)
# ──────────────────────────────────────────────────────────────────────────────


@st.cache_data(show_spinner=False)
def extract_pages_from_labels(labels_pdf_bytes: bytes, n_cols: int) -> List[PageData]:
    pages: List[PageData] = []
    with pdfplumber.open(io.BytesIO(labels_pdf_bytes)) as pdf:
        for pidx, page in enumerate(pdf.pages):
            words = page.extract_words(use_text_flow=True, keep_blank_chars=False) or []
            five = [
                w for w in words if re.fullmatch(r"\d{5}", str(w.get("text", "")))
            ]
            if not five:
                continue

            xs = [(w["x0"] + w["x1"]) / 2 for w in five]
            ys = [(w["top"] + w["bottom"]) / 2 for w in five]

            x_centers = kmeans_1d(xs, n_cols)
            y_centers = kmeans_1d(ys, max(1, round(len(five) / max(1, n_cols))))

            x_bounds = boundaries_from_centers(x_centers)
            y_bounds = boundaries_from_centers(y_centers)

            cell_map: Dict[Tuple[int, int], Tuple[float, dict]] = {}
            for w, xc, yc in zip(five, xs, ys):
                col = int(np.argmin(np.abs(x_centers - xc)))
                row = int(np.argmin(np.abs(y_centers - yc)))
                dist = abs(x_centers[col] - xc) + abs(y_centers[row] - yc)
                key = (row, col)
                if key not in cell_map or dist < cell_map[key][0]:
                    cell_map[key] = (dist, w)

            cells: List[CellData] = []
            for (row, col), (_, _w) in sorted(cell_map.items()):
                bbox = (
                    float(x_bounds[col]),
                    float(y_bounds[row]),
                    float(x_bounds[col + 1]),
                    float(y_bounds[row + 1]),
                )
                txt = (page.crop(bbox).extract_text() or "").strip()
                name, last5, qty = parse_label_cell_text(txt)
                cells.append(
                    CellData(
                        row=row,
                        col=col,
                        bbox=bbox,
                        name=name,
                        last5=last5,
                        qty=qty,
                        upc12=None,
                    )
                )

            pages.append(
                PageData(page_index=pidx, x_bounds=x_bounds, y_bounds=y_bounds, cells=cells)
            )
    return pages


# ──────────────────────────────────────────────────────────────────────────────
# Full-pallet page extraction
# ──────────────────────────────────────────────────────────────────────────────


def _wc(w: dict) -> Tuple[float, float]:
    return (w["x0"] + w["x1"]) / 2, (w["top"] + w["bottom"]) / 2


def _union(words: List[dict], px: float = 0, py: float = 0) -> Tuple[float, float, float, float]:
    return (
        min(w["x0"] for w in words) - px,
        min(w["top"] for w in words) - py,
        max(w["x1"] for w in words) + px,
        max(w["bottom"] for w in words) + py,
    )


def _group_nearby(words: List[dict], x_tol: float, y_tol: float) -> List[List[dict]]:
    groups: List[List[dict]] = []
    for w in sorted(words, key=lambda ww: (_wc(ww)[1], _wc(ww)[0])):
        cx, cy = _wc(w)
        placed = False
        for g in groups:
            bx0, bt, bx1, bb = _union(g)
            gcx, gcy = (bx0 + bx1) / 2, (bt + bb) / 2
            if abs(cx - gcx) <= max(x_tol, (bx1 - bx0) / 2 + x_tol * 0.4) and abs(
                cy - gcy
            ) <= max(y_tol, (bb - bt) / 2 + y_tol * 0.4):
                g.append(w)
                placed = True
                break
        if not placed:
            groups.append([w])

    changed = True
    while changed:
        changed = False
        merged: List[List[dict]] = []
        while groups:
            base = groups.pop(0)
            bx0, bt, bx1, bb = _union(base)
            i = 0
            while i < len(groups):
                gx0, gt, gx1, gb = _union(groups[i])
                if not (
                    bx1 + x_tol < gx0
                    or gx1 + x_tol < bx0
                    or bb + y_tol < gt
                    or gb + y_tol < bt
                ):
                    base.extend(groups.pop(i))
                    bx0, bt, bx1, bb = _union(base)
                    changed = True
                else:
                    i += 1
            merged.append(base)
        groups = merged

    return groups


@st.cache_data(show_spinner=False)
def extract_full_pallet_pages(labels_pdf_bytes: bytes) -> List[FullPalletPage]:
    pages: List[FullPalletPage] = []
    with pdfplumber.open(io.BytesIO(labels_pdf_bytes)) as pdf:
        for pidx, page in enumerate(pdf.pages):
            words = page.extract_words(use_text_flow=True, keep_blank_chars=False) or []
            five = [
                w for w in words if re.fullmatch(r"\d{5}", str(w.get("text", "")))
            ]
            if not five:
                continue

            pw, ph = float(page.width), float(page.height)

            # ── FIX B: exclude tokens in the gift-card-holder zone ───────────
            # The top ~31 % of each labels page contains the "WM GIFTCARD IN NEW
            # PKG" holder labels plus GCI packaging product codes (08470, 08478,
            # 08481, 08705, 08706, …).  These are physical display fixtures — not
            # the gift cards that belong in the main / bonus product grid.
            # Filtering them out prevents them from becoming phantom grid cells.
            holder_zone_bottom = ph * 0.31
            five = [w for w in five if float(w.get("top", 0)) >= holder_zone_bottom]
            if not five:
                continue

            xs = [(w["x0"] + w["x1"]) / 2 for w in five]
            ys = [(w["top"] + w["bottom"]) / 2 for w in five]

            # IMPORTANT: reduce over-splitting of columns
            x_centers = cluster_positions(xs, tol=max(10, pw * 0.025))
            y_centers = cluster_positions(ys, tol=max(7, ph * 0.012))
            if len(x_centers) == 0 or len(y_centers) == 0:
                continue

            x_bounds = boundaries_from_centers(x_centers)
            y_bounds = boundaries_from_centers(y_centers)
            if len(x_bounds) < 2 or len(y_bounds) < 2:
                continue

            cell_map: Dict[Tuple[int, int], Tuple[float, dict]] = {}
            for w, xc, yc in zip(five, xs, ys):
                col = int(np.argmin(np.abs(x_centers - xc)))
                row = int(np.argmin(np.abs(y_centers - yc)))
                dist = abs(x_centers[col] - xc) + abs(y_centers[row] - yc)
                key = (row, col)
                if key not in cell_map or dist < cell_map[key][0]:
                    cell_map[key] = (dist, w)

            cells: List[CellData] = []
            for (row, col), (_, token_w) in sorted(cell_map.items()):
                bbox = (
                    float(x_bounds[col]),
                    float(y_bounds[row]),
                    float(x_bounds[col + 1]),
                    float(y_bounds[row + 1]),
                )
                txt = (page.crop(bbox).extract_text() or "").strip()
                parsed_name, parsed_last5, qty = parse_label_cell_text(txt)
                token_last5 = str(token_w.get("text", "")).strip()
                raw_last5 = token_last5 if re.fullmatch(r"\d{5}", token_last5) else parsed_last5
                last5 = _to_last5(raw_last5)

                cells.append(
                    CellData(
                        row=row,
                        col=col,
                        bbox=bbox,
                        name=parsed_name,
                        last5=last5 or "",
                        qty=qty,
                        upc12=None,
                    )
                )

            annotations: List[AnnotationBox] = []
            wt = lambda w: str(w.get("text", "")).strip().upper()

            if xs:
                cx0_content = min(xs) - 15
                cx1_content = max(xs) + 15
            else:
                cx0_content, cx1_content = 150.0, 470.0

            wm_grp = [
                w
                for w in words
                if wt(w) in {"WM", "GIFTCAR", "GIFTCARD", "IN", "NEW", "PKG", "D"}
                and float(w["top"]) < ph * 0.30
            ]
            if wm_grp:
                by0 = min(w["top"] for w in wm_grp) - 8
                annotations.append(
                    AnnotationBox(
                        kind="gift_card_holders",
                        label="GIFT CARD HOLDERS",
                        bbox=(cx0_content, by0 - 22, cx1_content, by0 - 1),
                    )
                )

                sub_groups = _group_nearby(
                    [w for w in wm_grp if wt(w) in {"WM", "GIFTCAR", "GIFTCARD", "D", "IN", "NEW", "PKG"}],
                    x_tol=14,
                    y_tol=22,
                )
                for sg in sub_groups:
                    if len(sg) >= 3:
                        annotations.append(
                            AnnotationBox(
                                kind="wm_new_pkg",
                                label="WM GIFTCARD\nIN NEW PKG",
                                bbox=_union(sg, px=4, py=3),
                            )
                        )

            bonus_words = [w for w in words if wt(w) == "BONUS"]
            if bonus_words:
                bw = min(bonus_words, key=lambda w: float(w["top"]))
                bcy = (float(bw["top"]) + float(bw["bottom"])) / 2
                annotations.append(
                    AnnotationBox(
                        kind="bonus_strip",
                        label="BONUS",
                        bbox=(cx0_content, bcy - 12, cx1_content, bcy + 12),
                    )
                )

            mkt_words = [w for w in words if wt(w) in {"MARKETING", "MESSAGE", "PANEL"}]
            for grp in _group_nearby(mkt_words, x_tol=40, y_tol=20):
                if {wt(w) for w in grp} & {"MARKETING"}:
                    annotations.append(
                        AnnotationBox(
                            kind="marketing_signage",
                            label="MARKETING\nMESSAGE PANEL",
                            bbox=_union(grp, px=6, py=4),
                        )
                    )

            fraud_words = [w for w in words if wt(w) in {"FRAUD", "SIGNAGE"}]
            for grp in _group_nearby(fraud_words, x_tol=30, y_tol=16):
                if {wt(w) for w in grp} & {"FRAUD"}:
                    annotations.append(
                        AnnotationBox(
                            kind="fraud_signage",
                            label="FRAUD\nSIGNAGE",
                            bbox=_union(grp, px=6, py=4),
                        )
                    )

            pages.append(
                FullPalletPage(
                    page_index=pidx,
                    side_letter=chr(ord("A") + min(pidx, 3)),
                    cells=cells,
                    annotations=annotations,
                )
            )
    return pages


# ──────────────────────────────────────────────────────────────────────────────
# Card-image cropping helper
# ──────────────────────────────────────────────────────────────────────────────


def image_from_bytes(img_bytes: Optional[bytes]) -> Optional[Image.Image]:
    """Convert image bytes to PIL Image, or return None if bytes are None/empty."""
    if not img_bytes:
        return None
    try:
        return Image.open(io.BytesIO(img_bytes)).convert("RGBA")
    except Exception:
        return None


def crop_image_cell(
    images_doc: fitz.Document,
    page_index: int,
    bbox: Tuple[float, float, float, float],
    zoom: float = 2.6,
    inset: float = 0.045,
) -> Image.Image:
    page = images_doc.load_page(page_index)
    x0, top, x1, bottom = bbox
    w, h = x1 - x0, bottom - top
    rect = fitz.Rect(x0 + w * inset, top + h * inset, x1 - w * inset, bottom - h * inset)
    pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), clip=rect, alpha=False)
    return Image.frombytes("RGB", (pix.width, pix.height), pix.samples)


# ──────────────────────────────────────────────────────────────────────────────
# Drawing helpers
# ──────────────────────────────────────────────────────────────────────────────


def wrap_text(text: str, max_w: float, font: str, size: float) -> List[str]:
    parts = (text or "").split()
    lines: List[str] = []
    cur: List[str] = []
    for p in parts:
        trial = " ".join(cur + [p])
        if not cur or pdfmetrics.stringWidth(trial, font, size) <= max_w:
            cur.append(p)
        else:
            lines.append(" ".join(cur))
            cur = [p]
    if cur:
        lines.append(" ".join(cur))
    return lines


def _fit_font(
    text: str,
    font: str,
    max_w: float,
    max_h: float,
    lo: float,
    hi: float,
    step: float = 0.5,
) -> float:
    t = (text or "").strip()
    if not t:
        return lo
    s = hi
    while s >= lo:
        if s <= max_h and pdfmetrics.stringWidth(t, font, s) <= max_w:
            return s
        s -= step
    return lo


def _ellipsis(text: str, font: str, size: float, max_w: float) -> str:
    t = (text or "").strip()
    if not t:
        return ""
    if pdfmetrics.stringWidth(t, font, size) <= max_w:
        return t
    ell = "…"
    lo, hi = 0, len(t)
    best = ell
    while lo <= hi:
        mid = (lo + hi) // 2
        cand = t[:mid].rstrip() + ell
        if pdfmetrics.stringWidth(cand, font, size) <= max_w:
            best = cand
            lo = mid + 1
        else:
            hi = mid - 1
    return best


def _draw_gradient(
    c: canvas.Canvas,
    x: float,
    y: float,
    w: float,
    h: float,
    left: Tuple[float, float, float],
    right: Tuple[float, float, float],
    steps: int = 80,
) -> None:
    for i in range(max(8, steps)):
        t = i / (steps - 1)
        c.setFillColorRGB(
            left[0] * (1 - t) + right[0] * t,
            left[1] * (1 - t) + right[1] * t,
            left[2] * (1 - t) + right[2] * t,
        )
        xi = x + w * i / steps
        c.rect(xi, y, w / steps + 0.5, h, stroke=0, fill=1)


def _draw_header(
    c: canvas.Canvas,
    page_w: float,
    page_h: float,
    top_bar_h: float,
    title_text: str,
    right_label: str,
    logo_img: Optional[Image.Image],
    grad_l: Tuple[float, float, float],
    grad_r: Tuple[float, float, float],
) -> None:
    hy = page_h - top_bar_h
    _draw_gradient(c, 0, hy, page_w, top_bar_h, grad_l, grad_r, steps=120)

    tx = 14.0
    if logo_img:
        lw, lh = logo_img.size
        th = top_bar_h * 0.62
        r = th / max(1, lh)
        dw, dh = lw * r, lh * r
        c.drawImage(
            ImageReader(logo_img),
            14,
            hy + (top_bar_h - dh) / 2,
            dw,
            dh,
            mask="auto",
        )
        tx = 14 + dw + 10

    c.setFillColorRGB(1, 1, 1)
    fs = _fit_font(title_text, TITLE_FONT, page_w - tx - 130, top_bar_h * 0.72, 18, 36)
    c.setFont(TITLE_FONT, fs)
    c.drawString(tx, hy + (top_bar_h - fs) / 2 + 2, title_text)

    if right_label:
        rfs = _fit_font(right_label, TITLE_FONT, 120, top_bar_h * 0.68, 14, 28)
        rw = pdfmetrics.stringWidth(right_label, TITLE_FONT, rfs)
        c.setFont(TITLE_FONT, rfs)
        c.drawString(page_w - rw - 14, hy + (top_bar_h - rfs) / 2 + 2, right_label)

    c.setLineWidth(0.8)
    c.setStrokeColorRGB(0.88, 0.88, 0.92)
    c.line(0, hy, page_w, hy)


def _draw_cell_text_block(
    c: canvas.Canvas,
    x: float,
    y: float,
    w: float,
    h: float,
    name: str,
    upc12: Optional[str],
    last5: str,
    qty: Optional[int],
) -> None:
    px, py = 5, 4
    mw = max(12.0, w - px * 2)
    avail = max(10.0, h - py * 2)

    upc_str = upc12 if upc12 else (f"???????{last5}" if last5 else "")
    ns = 13.0
    while ns >= 5.0:
        ms = max(5.0, ns * 0.86)
        lhn = ns * 1.18
        lhm = ms * 1.18

        nls = wrap_text(name or "", mw, BODY_BOLD_FONT, ns)
        ul = f"UPC: {upc_str}" if upc_str else ""
        ql = f"Qty: {qty}" if qty is not None else ""

        need = (
            len(nls) * lhn
            + (2.5 if nls and (ul or ql) else 0)
            + (lhm if ul else 0)
            + (lhm if ql else 0)
        )
        if need <= avail:
            ty = y + h - py - ns
            c.setFillColorRGB(*NAVY_RGB)
            c.setFont(BODY_BOLD_FONT, ns)
            for ln in nls:
                c.drawString(x + px, max(y + py, ty), ln)
                ty -= lhn

            if nls and (ul or ql):
                ty -= 2.5
            c.setFont(BODY_FONT, ms)
            if ul:
                c.drawString(x + px, max(y + py, ty), ul)
                ty -= lhm
            if ql:
                c.drawString(x + px, max(y + py, ty), ql)
            return
        ns -= 0.5

    fs = 5.0
    ty = y + h - py - fs
    c.setFillColorRGB(*NAVY_RGB)
    c.setFont(BODY_BOLD_FONT, fs)
    for ln in wrap_text(name or "", mw, BODY_BOLD_FONT, fs):
        c.drawString(x + px, max(y + py, ty), ln)
        ty -= fs * 1.18
    c.setFont(BODY_FONT, fs)
    if upc_str:
        c.drawString(x + px, max(y + py, ty), f"UPC: {upc_str}")
        ty -= fs * 1.18
    if qty is not None:
        c.drawString(x + px, max(y + py, ty), f"Qty: {qty}")


_STOPWORDS = {
    "GIFT", "CARD", "CARDS", "VALUE", "THE", "AND", "FOR", "WITH", "IN",
    "OF", "A", "AN", "WALMART", "WM", "VGC", "GC",
}


def _is_important_token(tok: str) -> bool:
    t = tok.strip().upper()
    if not t:
        return False
    if "$" in t and re.search(r"\d", t):
        return True
    if re.fullmatch(r"\d{1,3}", t):
        return True
    if re.search(r"\d", t) and any(u in t for u in ("PK", "CT", "OZ", "LB", "MG", "ML")):
        return True
    if re.fullmatch(r"[A-Z]\d{1,3}", t):
        return True
    return False


def _compact_one_line_name(name: str) -> str:
    raw = re.sub(r"\s+", " ", (name or "").strip())
    if not raw:
        return ""
    tokens = raw.split()
    imp = [t for t in tokens if _is_important_token(t)]
    core = [t for t in tokens if t.upper() not in _STOPWORDS]
    use = core if core else tokens
    for t in imp:
        if t not in use:
            use.append(t)
    return " ".join(use)


def _fit_name_preserve_qualifiers(name: str, font: str, size: float, max_w: float) -> str:
    raw = re.sub(r"\s+", " ", (name or "").strip())
    if not raw:
        return ""
    tokens = raw.split()

    def is_important(t: str) -> bool:
        tt = t.upper()
        if "$" in tt and re.search(r"\d", tt):
            return True
        if re.search(r"\d", tt) and any(x in tt for x in ("PK", "CT", "OZ", "LB", "MG", "ML")):
            return True
        if re.fullmatch(r"[A-Z]\d{1,3}", tt):
            return True
        return False

    keep = tokens[:]
    while keep and pdfmetrics.stringWidth(" ".join(keep), font, size) > max_w:
        removable = [i for i, t in enumerate(keep) if 0 < i < len(keep) - 1 and not is_important(t)]
        if not removable:
            break
        mid = removable[len(removable) // 2]
        keep.pop(mid)

    candidate = " ".join(keep)
    if pdfmetrics.stringWidth(candidate, font, size) <= max_w:
        return candidate
    return _ellipsis(candidate, font, size, max_w)


def _draw_full_pallet_header(
    c: canvas.Canvas,
    page_w: float,
    page_h: float,
    header_h: float,
    title_text: str,
    right_label: str,
    logo_img: Optional[Image.Image],
) -> None:
    y0 = page_h - header_h
    c.setFillColorRGB(1, 1, 1)
    c.rect(0, y0, page_w, header_h, stroke=0, fill=1)

    tx = 18.0
    if logo_img:
        lw, lh = logo_img.size
        th = header_h * 0.62
        r = th / max(1, lh)
        dw, dh = lw * r, lh * r
        c.drawImage(
            ImageReader(logo_img),
            18,
            y0 + (header_h - dh) / 2,
            dw,
            dh,
            mask="auto",
        )
        tx = 18 + dw + 12

    c.setFillColorRGB(*NAVY_RGB)
    fs = _fit_font(title_text, BODY_BOLD_FONT, page_w - tx - 160, header_h * 0.70, 14, 26)
    c.setFont(BODY_BOLD_FONT, fs)
    c.drawString(tx, y0 + (header_h - fs) / 2 + 1, title_text)

    if right_label:
        rfs = _fit_font(right_label, BODY_BOLD_FONT, 150, header_h * 0.68, 12, 24)
        rw = pdfmetrics.stringWidth(right_label, BODY_BOLD_FONT, rfs)
        c.setFont(BODY_BOLD_FONT, rfs)
        c.drawString(page_w - rw - 18, y0 + (header_h - rfs) / 2 + 1, right_label)

    c.setLineWidth(0.8)
    c.setStrokeColorRGB(0.86, 0.86, 0.88)
    c.line(0, y0, page_w, y0)


def _draw_footer(c: canvas.Canvas, page_w: float, outer_margin: float, footer_h: float) -> None:
    fy = outer_margin + footer_h
    c.setLineWidth(0.7)
    c.setStrokeColorRGB(0.82, 0.82, 0.82)
    c.line(outer_margin, fy, page_w - outer_margin, fy)

    c.setFillColorRGB(*NAVY_RGB)
    c.setFont(BODY_BOLD_FONT, 11)
    left = f"Date: {date.today().isoformat()}"
    mid = "Generated by Kendal King"
    c.drawString(outer_margin, outer_margin + 11, left)
    mw = pdfmetrics.stringWidth(mid, BODY_BOLD_FONT, 11)
    c.drawString((page_w - mw) / 2, outer_margin + 11, mid)


# ──────────────────────────────────────────────────────────────────────────────
# Full-Pallet Display renderer
# ──────────────────────────────────────────────────────────────────────────────


def render_full_pallet_pdf(
    pages: List[FullPalletPage],
    images_pdf_bytes: bytes,
    matrix_idx: Dict[str, List[MatrixRow]],
    title_prefix: str = "POG",
    ppt_cards: Optional[Dict[str, PptSideCards]] = None,
    gift_holders: Optional[Dict[str, List[GiftHolder]]] = None,
    ppt_cpp_global: Optional[int] = None,
    debug: bool = False,
    debug_overlay: bool = False,
) -> bytes:
    buf = io.BytesIO()
    images_doc = fitz.open(stream=images_pdf_bytes, filetype="pdf")

    PAGE_W, BASE_PAGE_H = 792.0, 1224.0  # 11x17 portrait template; pages may grow taller
    MARGIN = 24.0
    HEADER_H = 56.0
    FOOTER_H = 36.0
    SECTION_BAR_H = 24.0
    SECTION_BAR_GAP = 10.0
    BUCKET_GAP = 12.0

    BASE_CONTENT_H = BASE_PAGE_H - (2 * MARGIN) - HEADER_H - FOOTER_H
    PPT_SECTION_H = BASE_CONTENT_H * 0.27
    HOLDER_SECTION_H = BASE_CONTENT_H * 0.18
    BAR_FILL = _hex_to_rgb("#77B5F0")
    BAR_TEXT = NAVY_RGB

    EMPTY_STROKE = (0.88, 0.88, 0.88)
    FILLED_STROKE = (0.78, 0.78, 0.80)

    logo = _try_load_logo()

    def _marker_y(p: FullPalletPage, kind: str) -> Optional[float]:
        for ann in p.annotations:
            if ann.kind == kind:
                _, t, _, b = ann.bbox
                return (t + b) / 2
        return None

    def _cell_center(cell: CellData) -> Tuple[float, float]:
        x0, t, x1, b = cell.bbox
        return (x0 + x1) / 2, (t + b) / 2

    def _global_row_order(p: FullPalletPage) -> List[int]:
        row_to_y: Dict[int, List[float]] = {}
        for cl in p.cells:
            _, y = _cell_center(cl)
            row_to_y.setdefault(cl.row, []).append(y)
        return sorted(row_to_y.keys(), key=lambda r: float(np.mean(row_to_y[r])))

    def _global_col_order_and_gaps(p: FullPalletPage) -> Tuple[List[int], List[int]]:
        col_to_x: Dict[int, List[float]] = {}
        for cl in p.cells:
            x, _ = _cell_center(cl)
            col_to_x.setdefault(cl.col, []).append(x)

        cols = sorted(col_to_x.keys(), key=lambda c_: float(np.mean(col_to_x[c_])))
        if len(cols) <= 1:
            return cols, []

        xs = [float(np.mean(col_to_x[c_])) for c_ in cols]
        diffs = [xs[i + 1] - xs[i] for i in range(len(xs) - 1)]
        med = float(np.median(diffs)) if diffs else 1.0
        if med <= 1e-6:
            return cols, [1] * (len(cols) - 1)

        units: List[int] = []
        for d in diffs:
            u = int(round(d / med))
            units.append(max(1, min(4, u)))
        return cols, units

    def _split_rows_by_bonus(p: FullPalletPage, global_rows: List[int]) -> Tuple[List[int], List[int]]:
        bonus_y = _marker_y(p, "bonus_strip")

        row_to_y: Dict[int, float] = {}
        for r in global_rows:
            ys = [(_cell_center(c)[1]) for c in p.cells if c.row == r]
            row_to_y[r] = float(np.mean(ys)) if ys else float("inf")

        yc = [row_to_y[r] for r in global_rows]

        if bonus_y is not None and len(yc) >= 2:
            above = [r for r in global_rows if row_to_y[r] < bonus_y]
            below = [r for r in global_rows if row_to_y[r] >= bonus_y]
            if above and below:
                return above, below

        if len(yc) < 2:
            return global_rows, []

        gaps = [yc[i + 1] - yc[i] for i in range(len(yc) - 1)]
        i_max = int(np.argmax(gaps))
        if gaps[i_max] < float(np.median(gaps)) * 1.6:
            return global_rows, []
        return global_rows[: i_max + 1], global_rows[i_max + 1 :]

    def _fit_x_layout(
        x0: float,
        x1: float,
        n_cols: int,
        gap_units: List[int],
        desired_card_w: float,
        desired_gap: float,
    ) -> Tuple[List[float], float, float, float, bool]:
        if n_cols <= 0:
            return [], desired_card_w, desired_gap, x0, False

        units = [max(1, int(gap_units[i])) if i < len(gap_units) else 1 for i in range(n_cols - 1)]
        avail_w = max(1.0, x1 - x0)
        gap = max(0.0, desired_gap)
        card_w = max(24.0, desired_card_w)

        def total_w(cw: float, gp: float) -> float:
            return n_cols * cw + sum(units) * gp

        t = total_w(card_w, gap)
        overflowed = t > avail_w

        if t > avail_w and sum(units) > 0:
            max_gap = (avail_w - n_cols * card_w) / sum(units)
            gap = max(2.0, min(gap, max_gap))
            t = total_w(card_w, gap)
        if t > avail_w:
            card_w = max(24.0, (avail_w - sum(units) * gap) / n_cols)
            t = total_w(card_w, gap)
        if t > avail_w:
            gap = 0.0
            card_w = max(24.0, avail_w / n_cols)
            t = total_w(card_w, gap)

        xs: List[float] = []
        cur = x0
        for i in range(n_cols):
            xs.append(cur)
            cur += card_w
            if i < n_cols - 1:
                cur += gap * units[i]

        right = xs[-1] + card_w
        if right > x1:
            shift = right - x1
            xs = [x - shift for x in xs]
            right = xs[-1] + card_w
            overflowed = True

        return xs, card_w, gap, right, overflowed

    def _section_cols_and_gaps(
        p: FullPalletPage,
        sec_rows: List[int],
        global_cols: List[int],
        global_gap_units: List[int],
    ) -> Tuple[List[int], List[int]]:
        if not sec_rows:
            return [], []
        row_set = set(sec_rows)
        global_rank = {c_: i for i, c_ in enumerate(global_cols)}
        sec_cols = sorted(
            {c.col for c in p.cells if c.row in row_set},
            key=lambda c_: global_rank.get(c_, 10**6 + c_),
        )
        if len(sec_cols) <= 1:
            return sec_cols, []
        gap_units: List[int] = []
        for i in range(len(sec_cols) - 1):
            lidx = global_rank.get(sec_cols[i], i)
            ridx = global_rank.get(sec_cols[i + 1], lidx + 1)
            if ridx <= lidx:
                gap_units.append(1)
                continue
            unit = sum(global_gap_units[lidx:ridx]) if global_gap_units else 1
            gap_units.append(max(1, int(unit)))
        return sec_cols, gap_units

    def _draw_section_bar(
        c: canvas.Canvas, x: float, y: float, w: float, h: float, label: str
    ) -> None:
        c.setFillColorRGB(*BAR_FILL)
        c.setStrokeColorRGB(*BAR_FILL)
        c.rect(x, y, w, h, stroke=0, fill=1)

        fs = _fit_font(label, BODY_BOLD_FONT, w - 20, h - 6, 10, 18, step=0.5)
        c.setFillColorRGB(*BAR_TEXT)
        c.setFont(BODY_BOLD_FONT, fs)
        tw = pdfmetrics.stringWidth(label, BODY_BOLD_FONT, fs)
        c.drawString(x + (w - tw) / 2, y + (h - fs) / 2 + 1, label)

    def _wrap_text_lines(
        text: str,
        font: str,
        size: float,
        max_w: float,
        max_lines: int,
    ) -> List[str]:
        words = [w for w in (text or "").strip().split() if w]
        if not words:
            return []

        lines_out: List[str] = []
        current_words: List[str] = []

        for word in words:
            trial_words = current_words + [word]
            trial = " ".join(trial_words)
            if not current_words or pdfmetrics.stringWidth(trial, font, size) <= max_w:
                current_words = trial_words
                continue

            lines_out.append(" ".join(current_words))
            current_words = [word]
            if len(lines_out) >= max_lines - 1:
                break

        if len(lines_out) < max_lines and current_words:
            last = " ".join(current_words)
            lines_out.append(_ellipsis(last, font, size, max_w))

        return lines_out[:max_lines]

    def _draw_card(
        c: canvas.Canvas,
        x: float,
        y: float,
        w: float,
        h: float,
        img: Optional[Image.Image],
        upc12: str,
        name: str,
        cpp: Optional[int],
    ) -> None:
        c.setFillColorRGB(1, 1, 1)
        c.setStrokeColorRGB(*FILLED_STROKE)
        c.setLineWidth(0.75)
        c.rect(x, y, w, h, stroke=1, fill=1)

        inner_pad_x = max(4.0, min(7.0, w * 0.045))
        inner_pad_y = max(4.0, min(6.0, h * 0.045))
        ix = x + inner_pad_x
        iy = y + inner_pad_y
        iw = max(8.0, w - 2 * inner_pad_x)
        ih = max(8.0, h - 2 * inner_pad_y)

        title_fs = 7.2 if w < 76 else 7.6 if w < 98 else 8.0
        meta_fs = 7.0 if w < 92 else 7.4

        footer_h = max(16.0, min(20.0, ih * 0.18))
        title_h = max(18.0, min(24.0, ih * 0.22))
        image_h = max(12.0, ih - footer_h - title_h - 6.0)

        footer_y = iy
        image_y = footer_y + footer_h + 2.0
        title_y = image_y + image_h + 4.0

        clip = c.beginPath()
        clip.rect(x + 0.6, y + 0.6, max(1.0, w - 1.2), max(1.0, h - 1.2))
        c.saveState()
        c.clipPath(clip, stroke=0, fill=0)

        if img is not None and iw > 8 and image_h > 8:
            sw, sh = img.size
            if sw > 0 and sh > 0:
                scale = min(iw / sw, image_h / sh)
                dw = sw * scale
                dh = sh * scale
                c.drawImage(
                    ImageReader(img),
                    ix + (iw - dw) / 2.0,
                    image_y + (image_h - dh) / 2.0,
                    dw,
                    dh,
                    preserveAspectRatio=True,
                    mask="auto",
                )

        title_text = _fit_name_preserve_qualifiers((name or "").upper(), BODY_FONT, title_fs, iw)
        title_lines = _wrap_text_lines(title_text, BODY_FONT, title_fs, iw, 2)

        c.setFillColorRGB(0.12, 0.12, 0.12)
        c.setFont(BODY_FONT, title_fs)
        if len(title_lines) == 1:
            tw = pdfmetrics.stringWidth(title_lines[0], BODY_FONT, title_fs)
            c.drawString(ix + max(0.0, (iw - tw) / 2.0), title_y + 5.0, title_lines[0])
        elif len(title_lines) >= 2:
            line_gap = 1.2
            line1_y = title_y + title_fs + line_gap
            line2_y = title_y
            for yy, line in ((line1_y, title_lines[0]), (line2_y, title_lines[1])):
                tw = pdfmetrics.stringWidth(line, BODY_FONT, title_fs)
                c.drawString(ix + max(0.0, (iw - tw) / 2.0), yy, line)

        upc_text = _ellipsis((upc12 or "").strip(), BODY_BOLD_FONT, meta_fs, iw)
        cpp_text = _ellipsis(
            f"CPP: {cpp}" if cpp is not None else "CPP:",
            BODY_BOLD_FONT,
            meta_fs,
            iw,
        )

        c.setFillColorRGB(0.05, 0.05, 0.05)
        c.setFont(BODY_BOLD_FONT, meta_fs)
        upc_tw = pdfmetrics.stringWidth(upc_text, BODY_BOLD_FONT, meta_fs)
        c.drawString(ix + max(0.0, (iw - upc_tw) / 2.0), footer_y + footer_h * 0.56, upc_text)

        cpp_tw = pdfmetrics.stringWidth(cpp_text, BODY_BOLD_FONT, meta_fs)
        c.drawString(ix + max(0.0, (iw - cpp_tw) / 2.0), footer_y + 1.5, cpp_text)

        c.restoreState()

    def _section_slot_policy(section_kind: str) -> Dict[str, float]:
        if section_kind == "main":
            return {
                "desired_card_w": 62.0,
                "min_card_w": 48.0,
                "max_card_w": 76.0,
                "card_ratio": 1.02,   # h / w
                "min_card_h": 58.0,
                "max_card_h": 86.0,
                "row_gutter": 10.0,
                "crop_zoom": 2.35,
                "crop_inset": 0.018,
            }

        return {
            "desired_card_w": 60.0,
            "min_card_w": 48.0,
            "max_card_w": 72.0,
            "card_ratio": 1.00,   # h / w
            "min_card_h": 56.0,
            "max_card_h": 82.0,
            "row_gutter": 10.0,
            "crop_zoom": 2.25,
            "crop_inset": 0.018,
        }

    def _measure_section_slot_map(
        p: FullPalletPage,
        sec_rows: List[int],
        section_kind: str,
        content_w: float,
        include_bar: bool,
    ) -> Dict[str, object]:
        if not sec_rows:
            return {
                "rows": [],
                "card_w": 0.0,
                "card_h": 0.0,
                "row_gutter": 0.0,
                "crop_zoom": 2.25,
                "crop_inset": 0.018,
                "total_h": 0.0,
                "leftmost": 0.0,
                "rightmost": 0.0,
                "overflow": False,
            }

        policy = _section_slot_policy(section_kind)
        row_set = set(sec_rows)

        rows_data: List[dict] = []
        all_cells: List[CellData] = []
        all_centers_src: List[float] = []
        min_gap_src: Optional[float] = None

        for r in sec_rows:
            row_cells = [cell for cell in p.cells if cell.row == r and cell.row in row_set]
            row_cells.sort(key=lambda cell: ((cell.bbox[0] + cell.bbox[2]) / 2.0, cell.col))

            if not row_cells:
                continue

            row_centers = [((cell.bbox[0] + cell.bbox[2]) / 2.0) for cell in row_cells]
            for i in range(1, len(row_centers)):
                gap = row_centers[i] - row_centers[i - 1]
                if gap > 0:
                    min_gap_src = gap if min_gap_src is None else min(min_gap_src, gap)

            rows_data.append(
                {
                    "row_id": r,
                    "cells": row_cells,
                    "centers_src": row_centers,
                }
            )
            all_cells.extend(row_cells)
            all_centers_src.extend(row_centers)

        if not rows_data or not all_cells or not all_centers_src:
            return {
                "rows": [],
                "card_w": 0.0,
                "card_h": 0.0,
                "row_gutter": 0.0,
                "crop_zoom": policy["crop_zoom"],
                "crop_inset": policy["crop_inset"],
                "total_h": 0.0,
                "leftmost": 0.0,
                "rightmost": 0.0,
                "overflow": False,
            }

        src_left = min(all_centers_src)
        src_right = max(all_centers_src)
        src_span = max(1.0, src_right - src_left)

        if min_gap_src is None:
            scaled_gap = policy["desired_card_w"] + 10.0
        else:
            center_scale = max(1.0, (content_w - 24.0) / src_span)
            scaled_gap = min_gap_src * center_scale

        card_w = min(
            policy["max_card_w"],
            max(policy["min_card_w"], min(policy["desired_card_w"], scaled_gap * 0.76)),
        )
        card_h = min(
            policy["max_card_h"],
            max(policy["min_card_h"], card_w * policy["card_ratio"]),
        )

        side_pad = (card_w / 2.0) + 6.0
        usable_w = max(40.0, content_w - (2.0 * side_pad))

        leftmost = float("inf")
        rightmost = float("-inf")

        for row_info in rows_data:
            centers_norm: List[float] = []
            for cx_src in row_info["centers_src"]:
                nx = side_pad + ((cx_src - src_left) / src_span) * usable_w
                centers_norm.append(nx)
                leftmost = min(leftmost, nx - card_w / 2.0)
                rightmost = max(rightmost, nx + card_w / 2.0)
            row_info["centers_norm"] = centers_norm

        row_gutter = policy["row_gutter"]
        total_h = len(rows_data) * card_h + max(0, len(rows_data) - 1) * row_gutter
        if include_bar:
            total_h += SECTION_BAR_H + SECTION_BAR_GAP

        overflow = leftmost < -0.001 or rightmost > content_w + 0.001

        return {
            "rows": rows_data,
            "card_w": card_w,
            "card_h": card_h,
            "row_gutter": row_gutter,
            "crop_zoom": policy["crop_zoom"],
            "crop_inset": policy["crop_inset"],
            "total_h": total_h,
            "leftmost": leftmost,
            "rightmost": rightmost,
            "overflow": overflow,
        }

    def _draw_slot_map_section(
        p: FullPalletPage,
        plan: Dict[str, object],
        sec_top: float,
        label: Optional[str],
        unresolved_bucket: List[str],
        content_x0: float,
        content_w: float,
        product_map: Dict[Tuple[int, int], Tuple[Optional[MatrixRow], CellData]],
    ) -> Tuple[int, int, bool, float]:
        nonlocal rightmost_used, matched_cells, unmatched_cells

        rows_data = plan["rows"]
        if not rows_data:
            return 0, 0, False, sec_top

        y_cursor = sec_top
        if label is not None:
            bar_y = y_cursor - SECTION_BAR_H
            _draw_section_bar(c, content_x0, bar_y, content_w, SECTION_BAR_H, label)
            y_cursor = bar_y - SECTION_BAR_GAP

        card_w = float(plan["card_w"])
        card_h = float(plan["card_h"])
        row_gutter = float(plan["row_gutter"])

        lowest_bottom = y_cursor
        max_cols = 0

        for ri, row_info in enumerate(rows_data):
            row_cells: List[CellData] = row_info["cells"]
            row_centers: List[float] = row_info["centers_norm"]
            max_cols = max(max_cols, len(row_cells))

            y = y_cursor - (ri + 1) * card_h - ri * row_gutter

            for cell, center_x in zip(row_cells, row_centers):
                x = content_x0 + float(center_x) - (card_w / 2.0)
                key = (cell.row, cell.col)
                match, _cell = product_map.get(key, (None, cell))

                last5_key = _to_last5(cell.last5)
                upc12 = match.upc12 if match else None
                cpp = match.cpp_qty if match else None
                disp_name = (match.display_name if match and match.display_name else cell.name).strip()

                if upc12:
                    upc_str = upc12
                else:
                    upc_str = f"LAST5 {last5_key}"
                    disp_name = f"UNRESOLVED {last5_key}"
                    unresolved_bucket.append(last5_key)

                try:
                    img = crop_image_cell(
                        images_doc,
                        p.page_index,
                        cell.bbox,
                        zoom=float(plan["crop_zoom"]),
                        inset=float(plan["crop_inset"]),
                    )
                except Exception:
                    img = None

                if match:
                    matched_cells += 1
                else:
                    unmatched_cells += 1

                _draw_card(c, x, y, card_w, card_h, img, upc_str, disp_name, cpp)
                rightmost_used = max(rightmost_used, x + card_w)

            lowest_bottom = min(lowest_bottom, y)

        return len(rows_data), max_cols, bool(plan["overflow"]), lowest_bottom

    def _draw_debug_box(
        c: canvas.Canvas, x: float, y: float, w: float, h: float, label: str
    ) -> None:
        c.setStrokeColorRGB(0.92, 0.20, 0.20)
        c.setLineWidth(0.7)
        c.rect(x, y, w, h, stroke=1, fill=0)
        c.setFillColorRGB(0.75, 0.10, 0.10)
        c.setFont("Helvetica", 7)
        c.drawString(x + 2, y + h - 9, label)

    try:
        if not pages:
            c = canvas.Canvas(buf, pagesize=(PAGE_W, BASE_PAGE_H))
            c.save()
            return buf.getvalue()

        c = canvas.Canvas(buf, pagesize=(PAGE_W, BASE_PAGE_H))

        for pdata in pages:
            cx0, cx1 = MARGIN, PAGE_W - MARGIN
            content_w = cx1 - cx0

            global_cols, global_gap_units = _global_col_order_and_gaps(pdata)
            global_rows = _global_row_order(pdata)
            above_bonus_rows, below_bonus_rows = _split_rows_by_bonus(pdata, global_rows)

            side_ppt = (
                ppt_cards.get(pdata.side_letter, PptSideCards(pdata.side_letter, [], []))
                if ppt_cards
                else PptSideCards(pdata.side_letter, [], [])
            )
            side_holders = gift_holders.get(pdata.side_letter, []) if gift_holders else []

            product_map: Dict[Tuple[int, int], Tuple[Optional[MatrixRow], CellData]] = {}
            for cell in pdata.cells:
                if cell.last5:
                    match = resolve_full_pallet(cell.last5, cell.name, matrix_idx)
                    product_map[(cell.row, cell.col)] = (match, cell)

            main_plan = _measure_section_slot_map(
                pdata,
                above_bonus_rows,
                "main",
                content_w,
                include_bar=False,
            )
            bonus_plan = _measure_section_slot_map(
                pdata,
                below_bonus_rows,
                "bonus",
                content_w,
                include_bar=bool(below_bonus_rows),
            )

            products_block_h = float(main_plan["total_h"])
            if above_bonus_rows and below_bonus_rows:
                products_block_h += BUCKET_GAP
            products_block_h += float(bonus_plan["total_h"])

            required_content_h = (
                PPT_SECTION_H
                + BUCKET_GAP
                + HOLDER_SECTION_H
                + BUCKET_GAP
                + products_block_h
            )
            PAGE_H = max(BASE_PAGE_H, (2 * MARGIN) + HEADER_H + FOOTER_H + required_content_h)

            c.setPageSize((PAGE_W, PAGE_H))

            _draw_full_pallet_header(
                c,
                PAGE_W,
                PAGE_H,
                HEADER_H,
                title_prefix.strip() or "POG",
                f"SIDE {pdata.side_letter}",
                logo,
            )
            _draw_footer(c, PAGE_W, MARGIN, FOOTER_H)

            cy0, cy1 = MARGIN + FOOTER_H, PAGE_H - MARGIN - HEADER_H
            content_h = cy1 - cy0

            bucket_a_h = PPT_SECTION_H
            bucket_b_h = HOLDER_SECTION_H

            bucket_a_top = cy1
            bucket_a_bottom = bucket_a_top - bucket_a_h
            bucket_b_top = bucket_a_bottom - BUCKET_GAP
            bucket_b_bottom = bucket_b_top - bucket_b_h
            products_top = bucket_b_bottom - BUCKET_GAP
            products_bottom = cy0

            if debug_overlay:
                c.setStrokeColorRGB(0.20, 0.70, 0.25)
                c.setLineWidth(0.8)
                c.rect(cx0, cy0, cx1 - cx0, cy1 - cy0, stroke=1, fill=0)
                c.setStrokeColorRGB(0.95, 0.10, 0.10)
                c.line(cx1, cy0, cx1, cy1)

            rightmost_used = cx0
            adjusted_to_fit = False
            matched_cells = 0
            unmatched_cells = 0
            unresolved_main: List[str] = []
            unresolved_bonus: List[str] = []

            # Bucket A: PPT Top Cards
            a_gap_y = 8.0
            top_card_h = max(30.0, min(86.0, bucket_a_h * 0.34))
            # Desired card width spans full content width (8 equal slots with gaps).
            _desired_top_card_w = max(72.0, (content_w - 7 * 6.0) / 8)
            top_xs, top_card_w, _, top_right, top_overflow = _fit_x_layout(
                cx0, cx1, 8, [1] * 7, _desired_top_card_w, 6.0
            )
            adjusted_to_fit = adjusted_to_fit or top_overflow
            rightmost_used = max(rightmost_used, top_right)

            top_row_y = bucket_a_top - top_card_h
            for i in range(8):
                x = top_xs[i]
                card = side_ppt.top8[i] if i < len(side_ppt.top8) else None
                c.setFillColorRGB(1, 1, 1)
                c.setStrokeColorRGB(*FILLED_STROKE if card else EMPTY_STROKE)
                c.setLineWidth(0.75 if card else 0.45)
                c.rect(x, top_row_y, top_card_w, top_card_h, stroke=1, fill=1)
                if card:
                    _draw_card(
                        c, x, top_row_y, top_card_w, top_card_h,
                        image_from_bytes(card.image_bytes),
                        f"ID# {card.card_id}", card.title, ppt_cpp_global,
                    )

            # Side panel: 3 columns × 2 rows, right-aligned.
            # Correct formula: block must be wide enough for 3 cards + 2 gaps.
            _side_gap = 6.0
            _side_cols = 3
            side_block_w = _side_cols * top_card_w + (_side_cols - 1) * _side_gap
            sx0 = max(cx0, cx1 - side_block_w)
            side_xs = [sx0 + i * (top_card_w + _side_gap) for i in range(_side_cols)]
            side_top = top_row_y - SECTION_BAR_GAP
            side_available_h = max(36.0, side_top - bucket_a_bottom)
            side_card_h = max(20.0, min(top_card_h, (side_available_h - a_gap_y) / 2))
            # 3 cols × 2 rows, row-major: indices 0-5
            for row in range(2):
                y = side_top - (row + 1) * side_card_h - row * a_gap_y
                for col in range(_side_cols):
                    idx = row * _side_cols + col
                    x = side_xs[col]
                    card = side_ppt.side6[idx] if idx < len(side_ppt.side6) else None
                    c.setFillColorRGB(1, 1, 1)
                    c.setStrokeColorRGB(*FILLED_STROKE if card else EMPTY_STROKE)
                    c.setLineWidth(0.75 if card else 0.45)
                    c.rect(x, y, top_card_w, side_card_h, stroke=1, fill=1)
                    if card:
                        _draw_card(
                            c, x, y, top_card_w, side_card_h,
                            image_from_bytes(card.image_bytes),
                            f"ID# {card.card_id}", card.title, ppt_cpp_global,
                        )
                    rightmost_used = max(rightmost_used, x + top_card_w)

            if debug_overlay:
                _draw_debug_box(c, cx0, bucket_a_bottom, content_w, bucket_a_top - bucket_a_bottom, "PPT")

            # Bucket B: Gift Card Holders
            holder_bar_y = bucket_b_top - SECTION_BAR_H
            _draw_section_bar(c, cx0, holder_bar_y, content_w, SECTION_BAR_H, "GIFT CARD HOLDERS")
            holder_grid_top = holder_bar_y - SECTION_BAR_GAP
            holder_grid_h = max(24.0, holder_grid_top - bucket_b_bottom)
            holder_cols = max(1, min(8, len(global_cols) if global_cols else 8))
            holder_rows = max(1, int(math.ceil(len(side_holders) / holder_cols)))
            holder_gutter_y = 6.0
            holder_card_h = max(16.0, (holder_grid_h - max(0, holder_rows - 1) * holder_gutter_y) / holder_rows)

            holder_xs, holder_card_w, _, holder_right, holder_overflow = _fit_x_layout(
                cx0, cx1, holder_cols, [1] * max(0, holder_cols - 1), 74.0, 6.0
            )
            adjusted_to_fit = adjusted_to_fit or holder_overflow
            rightmost_used = max(rightmost_used, holder_right)

            gift_top = holder_grid_top
            for ri in range(holder_rows):
                y = gift_top - (ri + 1) * holder_card_h - ri * holder_gutter_y
                for ci in range(holder_cols):
                    idx = ri * holder_cols + ci
                    x = holder_xs[ci]
                    holder = side_holders[idx] if idx < len(side_holders) else None
                    c.setFillColorRGB(1, 1, 1)
                    c.setStrokeColorRGB(*FILLED_STROKE if holder else EMPTY_STROKE)
                    c.setLineWidth(0.75 if holder else 0.45)
                    c.rect(x, y, holder_card_w, holder_card_h, stroke=1, fill=1)
                    if holder:
                        _draw_card(
                            c,
                            x,
                            y,
                            holder_card_w,
                            holder_card_h,
                            image_from_bytes(holder.image_bytes),
                            holder.item_no,
                            holder.name,
                            holder.qty,
                        )

            if debug_overlay:
                _draw_debug_box(c, cx0, bucket_b_bottom, content_w, bucket_b_top - bucket_b_bottom, "HOLDERS")

            main_rows_count = 0
            main_cols = 0
            main_over = False
            main_bottom = products_top

            if above_bonus_rows:
                (
                    main_rows_count,
                    main_cols,
                    main_over,
                    main_bottom,
                ) = _draw_slot_map_section(
                    pdata,
                    main_plan,
                    products_top,
                    None,
                    unresolved_main,
                    cx0,
                    content_w,
                    product_map,
                )

            bonus_top = main_bottom - (BUCKET_GAP if above_bonus_rows and below_bonus_rows else 0.0)
            bonus_rows_count = 0
            bonus_cols = 0
            bonus_over = False
            bonus_bottom = bonus_top

            if below_bonus_rows:
                (
                    bonus_rows_count,
                    bonus_cols,
                    bonus_over,
                    bonus_bottom,
                ) = _draw_slot_map_section(
                    pdata,
                    bonus_plan,
                    bonus_top,
                    "BONUS",
                    unresolved_bonus,
                    cx0,
                    content_w,
                    product_map,
                )

            adjusted_to_fit = adjusted_to_fit or main_over or bonus_over

            if debug_overlay:
                _draw_debug_box(c, cx0, products_bottom, content_w, products_top - products_bottom, "PRODUCTS")
                if above_bonus_rows:
                    _draw_debug_box(
                        c,
                        cx0,
                        main_bottom,
                        content_w,
                        products_top - main_bottom,
                        "MAIN SLOT MAP",
                    )
                if below_bonus_rows:
                    _draw_debug_box(
                        c,
                        cx0,
                        bonus_bottom,
                        content_w,
                        bonus_top - bonus_bottom,
                        "BONUS SLOT MAP",
                    )

            right_limit = cx1
            exceeded = rightmost_used > right_limit + 0.001
            lowest_bottom = bonus_bottom if below_bonus_rows else main_bottom
            vertical_overflow = lowest_bottom < products_bottom - 0.001
            adjusted_to_fit = adjusted_to_fit or exceeded or vertical_overflow

            main_last5_codes = sorted({_to_last5(c.last5) for c in pdata.cells if c.row in set(above_bonus_rows)})
            bonus_last5_codes = sorted({_to_last5(c.last5) for c in pdata.cells if c.row in set(below_bonus_rows)})
            main_slots_total = sum(1 for c in pdata.cells if c.row in set(above_bonus_rows))
            bonus_slots_total = sum(1 for c in pdata.cells if c.row in set(below_bonus_rows))

            if debug:
                st.write(
                    {
                        "side": pdata.side_letter,
                        "page_size": f"{PAGE_W}x{PAGE_H}",
                        "margins": {"left": cx0, "right": cx1, "top": cy1, "bottom": cy0},
                        "bucket_regions": {
                            "ppt": [round(bucket_a_top, 1), round(bucket_a_bottom, 1)],
                            "holders": [round(bucket_b_top, 1), round(bucket_b_bottom, 1)],
                            "products_available": [round(products_top, 1), round(products_bottom, 1)],
                            "main_used": [round(products_top, 1), round(main_bottom, 1)] if above_bonus_rows else None,
                            "bonus_used": [round(bonus_top, 1), round(bonus_bottom, 1)] if below_bonus_rows else None,
                        },
                        "grid_detected": {
                            "main_cols": main_cols,
                            "main_rows": main_rows_count,
                            "bonus_cols": bonus_cols,
                            "bonus_rows": bonus_rows_count,
                        },
                        "layout_width": {
                            "rightmost_x": round(rightmost_used, 2),
                            "right_limit": round(right_limit, 2),
                            "exceeded": exceeded,
                            "vertical_overflow": vertical_overflow,
                            "adjusted_to_fit": adjusted_to_fit,
                        },
                        "ppt_counts": {
                            "top8": len(side_ppt.top8),
                            "side6": len(side_ppt.side6),
                            "ids_top8": [c_.card_id for c_ in side_ppt.top8],
                            "ids_side6": [c_.card_id for c_ in side_ppt.side6],
                        },
                        "holders": {
                            "count": len(side_holders),
                            "items": [h.item_no for h in side_holders],
                        },
                        "slot_codes": {
                            "main_total_slots": main_slots_total,
                            "bonus_total_slots": bonus_slots_total,
                            "main_unique_count": len(main_last5_codes),
                            "bonus_unique_count": len(bonus_last5_codes),
                            "main": main_last5_codes,
                            "bonus": bonus_last5_codes,
                        },
                        "unresolved_last5": {
                            "main": sorted({x for x in unresolved_main if x}),
                            "bonus": sorted({x for x in unresolved_bonus if x}),
                        },
                        "matrix_matches": {
                            "matched": matched_cells,
                            "unmatched": unmatched_cells,
                        },
                    }
                )

            c.showPage()

        c.save()
        return buf.getvalue()
    finally:
        images_doc.close()


# ──────────────────────────────────────────────────────────────────────────────
# Standard Flat Display renderer (UNCHANGED)
# ──────────────────────────────────────────────────────────────────────────────


def render_standard_pog_pdf(
    pages: List[PageData],
    images_pdf_bytes: bytes,
    matrix_idx: Dict[str, List[MatrixRow]],
    title_prefix: str = "POG",
) -> bytes:
    buf = io.BytesIO()
    images_doc = fitz.open(stream=images_pdf_bytes, filetype="pdf")

    scale_factor = 1.5
    outer_margin = 44.0
    side_gap = 28.0
    top_bar_h = 90.0
    footer_h = 44.0
    side_label_h = 56.0
    cell_inset = 5.0
    border_w = 0.75
    img_frac = 0.58

    grad_left = _hex_to_rgb("#5B63A9")
    grad_right = _hex_to_rgb("#3E4577")
    logo_img = _try_load_logo()

    side_count = len(pages)
    per_side_w = int(310 * scale_factor)

    side_scales: List[float] = []
    side_heights: List[float] = []
    for page in pages:
        x_min = float(page.x_bounds[0])
        x_max = float(page.x_bounds[-1])
        y_min = float(page.y_bounds[0])
        y_max = float(page.y_bounds[-1])

        scale = per_side_w / max(1e-6, x_max - x_min)
        side_scales.append(scale)
        side_heights.append(scale * max(1e-6, y_max - y_min))

    content_h = max(side_heights) if side_heights else 600.0
    page_w = outer_margin * 2 + side_count * per_side_w + max(0, side_count - 1) * side_gap
    page_h = outer_margin + top_bar_h + side_label_h + content_h + footer_h + outer_margin

    c = canvas.Canvas(buf, pagesize=(page_w, page_h))

    try:
        _draw_header(
            c, page_w, page_h, top_bar_h,
            title_prefix or "POG", "", logo_img, grad_left, grad_right,
        )

        cells_top = page_h - top_bar_h - side_label_h
        content_bottom = outer_margin + footer_h

        _draw_footer(c, page_w, outer_margin, footer_h)

        for side_idx, page in enumerate(pages):
            side_letter = chr(ord("A") + side_idx)
            side_origin_x = outer_margin + side_idx * (per_side_w + side_gap)

            badge_h = 34.0
            badge_w = 148.0
            badge_y = cells_top + (side_label_h - badge_h) / 2

            c.setFillColorRGB(1, 1, 1)
            c.setStrokeColorRGB(0.85, 0.85, 0.90)
            c.setLineWidth(0.85)
            c.roundRect(side_origin_x, badge_y, badge_w, badge_h, 8, stroke=1, fill=1)

            side_text = f"Side {side_letter}"
            side_font_size = _fit_font(side_text, TITLE_FONT, badge_w - 16, badge_h - 8, 14, 22)
            side_text_w = pdfmetrics.stringWidth(side_text, TITLE_FONT, side_font_size)
            c.setFillColorRGB(*NAVY_RGB)
            c.setFont(TITLE_FONT, side_font_size)
            c.drawString(
                side_origin_x + (badge_w - side_text_w) / 2,
                badge_y + (badge_h - side_font_size) / 2 + 1,
                side_text,
            )

            if side_idx > 0:
                sep_x = side_origin_x - side_gap / 2
                c.setLineWidth(0.6)
                c.setStrokeColorRGB(0.88, 0.88, 0.88)
                c.line(sep_x, content_bottom + 2, sep_x, page_h - top_bar_h)

            x_min = float(page.x_bounds[0])
            y_min = float(page.y_bounds[0])
            scale = side_scales[side_idx]

            for cell in page.cells:
                match = _resolve(cell.last5, cell.name, matrix_idx) if cell.last5 else None
                upc12 = match.upc12 if match else None
                qty = cell.qty if cell.qty is not None else (match.cpp_qty if match else None)

                x0, top, x1, bottom = cell.bbox
                out_x0 = side_origin_x + (x0 - x_min) * scale + cell_inset
                out_x1 = side_origin_x + (x1 - x_min) * scale - cell_inset
                out_top = cells_top - (top - y_min) * scale - cell_inset
                out_bottom = cells_top - (bottom - y_min) * scale + cell_inset

                out_w = out_x1 - out_x0
                out_h = out_top - out_bottom
                if out_w <= 2 or out_h <= 2:
                    continue

                c.setFillColorRGB(1, 1, 1)
                c.setStrokeColorRGB(0.72, 0.72, 0.72)
                c.setLineWidth(border_w)
                c.rect(out_x0, out_bottom, out_w, out_h, stroke=1, fill=1)

                img_area_h = out_h * img_frac
                text_area_h = out_h - img_area_h

                img = crop_image_cell(images_doc, page.page_index, cell.bbox, zoom=3.2, inset=0.08)
                img_w, img_h = img.size
                img_scale = min(out_w * 0.86 / max(1, img_w), img_area_h * 0.84 / max(1, img_h))
                draw_w, draw_h = img_w * img_scale, img_h * img_scale

                c.drawImage(
                    ImageReader(img),
                    out_x0 + (out_w - draw_w) / 2,
                    out_bottom + text_area_h + (img_area_h - draw_h) / 2,
                    draw_w,
                    draw_h,
                    preserveAspectRatio=True,
                    mask="auto",
                )

                c.setLineWidth(0.5)
                c.setStrokeColorRGB(0.88, 0.88, 0.88)
                c.line(out_x0 + 3, out_bottom + text_area_h, out_x1 - 3, out_bottom + text_area_h)

                _draw_cell_text_block(
                    c, out_x0, out_bottom, out_w, text_area_h,
                    cell.name, upc12, cell.last5, qty,
                )

        c.showPage()
        c.save()
        return buf.getvalue()
    finally:
        images_doc.close()


# ──────────────────────────────────────────────────────────────────────────────
# Streamlit UI
# ──────────────────────────────────────────────────────────────────────────────


def main() -> None:
    st.set_page_config(page_title="Planogram Generator", layout="wide")
    st.title("Planogram Generator")

    with st.sidebar:
        st.header("Configuration")
        display_type = st.selectbox(
            "Display type",
            [DISPLAY_STANDARD, DISPLAY_FULL_PALLET],
            index=0,
        )

        st.divider()
        matrix_file = st.file_uploader("Matrix Excel (.xlsx)", type=["xlsx"])
        labels_pdf = st.file_uploader("Labels PDF", type=["pdf"])
        images_pdf = st.file_uploader("Images PDF", type=["pdf"])

        st.divider()
        title_prefix = st.text_input("PDF title prefix", "POG")
        out_name = st.text_input("Output filename", "pog_export.pdf")
        generate = st.button("Generate POG PDF", type="primary", use_container_width=True)

    if not (matrix_file and labels_pdf and images_pdf):
        st.info("Upload Matrix XLSX + Labels PDF + Images PDF to begin.")
        return

    matrix_bytes = matrix_file.getvalue()
    labels_bytes = labels_pdf.getvalue()
    images_bytes = images_pdf.getvalue()

    matrix_idx = load_matrix_index(matrix_bytes)

    if display_type == DISPLAY_STANDARD:
        pages = extract_pages_from_labels(labels_bytes, N_COLS)
        if not pages:
            st.error("No 5-digit UPC tokens found in Labels PDF.")
            return

        st.subheader(f"Detected {len(pages)} side(s)")
        rows = []
        for i, p in enumerate(pages):
            sl = chr(ord("A") + i)
            for cell in p.cells:
                match = _resolve(cell.last5, cell.name, matrix_idx) if cell.last5 else None
                rows.append(
                    {
                        "Side": sl,
                        "Row": cell.row,
                        "Col": cell.col,
                        "Name": cell.name,
                        "Last5": cell.last5,
                        "Qty": cell.qty if cell.qty is not None else (match.cpp_qty if match else None),
                        "UPC12": match.upc12 if match else None,
                    }
                )

        st.dataframe(
            pd.DataFrame(rows).sort_values(["Side", "Row", "Col"]),
            use_container_width=True,
            height=420,
        )

        if generate:
            with st.spinner("Rendering PDF…"):
                pdf = render_standard_pog_pdf(
                    pages, images_bytes, matrix_idx, title_prefix.strip() or "POG",
                )
            st.success("Done.")
            st.download_button(
                "⬇ Download Planogram PDF",
                pdf,
                file_name=out_name if out_name.endswith(".pdf") else f"{out_name}.pdf",
                mime="application/pdf",
                use_container_width=True,
            )

    else:
        pptx_file = st.file_uploader("Top Cards Blueprint (.pptx)", type=["pptx"])
        gift_file = st.file_uploader("2025 D82 POG Workbook (.xlsx)", type=["xlsx"])
        ppt_cpp_global = st.number_input("PPT Cards CPP (Global)", min_value=0, value=0, step=1)
        show_debug = st.checkbox("Show debug details")
        show_layout_overlay = st.checkbox("Show Full Pallet layout overlay")

        if not (pptx_file and gift_file):
            st.info("Upload PPTX + POG XLSX for Full Pallet mode.")
            return

        try:
            ppt_cards = load_ppt_cards(pptx_file.getvalue())
        except ImportError:
            st.error(
                "python-pptx is not installed. Full Pallet mode requires python-pptx to parse the Top Cards Blueprint."
            )
            return
        except Exception as e:
            if show_debug:
                st.exception(e)
            st.error("Unable to parse Top Cards Blueprint (.pptx). Please verify the file.")
            return

        ppt_issues = validate_ppt_side_cards(ppt_cards)
        if ppt_issues:
            st.error("Top Cards Blueprint validation failed:\n" + "\n".join(f"- {msg}" for msg in ppt_issues))
            return

        if ppt_cpp_global in (None, 0):
            st.warning("PPT CPP global is 0 or empty. PPT cards will render with CPP: 0.")

        try:
            gift_holders = load_gift_card_holders(gift_file.getvalue())
        except Exception as e:
            if show_debug:
                st.exception(e)
            st.error(f"Unable to parse POG workbook holder table: {e}")
            return

        try:
            fp_matrix_idx = load_full_pallet_matrix_index(matrix_bytes)
        except Exception as e:
            if show_debug:
                st.exception(e)
            st.error("Unable to parse D82 Item List for Full Pallet matching.")
            return

        fp_pages = extract_full_pallet_pages(labels_bytes)
        if not fp_pages:
            st.error("No product cells detected in Labels PDF.")
            return

        st.subheader(f"Detected {len(fp_pages)} side(s) — one output page per side")
        rows = []
        for pg in fp_pages:
            for cell in pg.cells:
                match = resolve_full_pallet(cell.last5, cell.name, fp_matrix_idx) if cell.last5 else None
                rows.append(
                    {
                        "Side": pg.side_letter,
                        "Row": cell.row,
                        "Col": cell.col,
                        "Name": cell.name,
                        "Last5": cell.last5,
                        "CPP (Excel)": match.cpp_qty if match else None,
                        "UPC12": match.upc12 if match else None,
                    }
                )

        st.dataframe(
            pd.DataFrame(rows).sort_values(["Side", "Row", "Col"]),
            use_container_width=True,
            height=420,
        )

        if generate:
            with st.spinner("Rendering PDF…"):
                try:
                    pdf = render_full_pallet_pdf(
                        fp_pages,
                        images_bytes,
                        fp_matrix_idx,
                        title_prefix.strip() or "POG",
                        ppt_cards=ppt_cards,
                        gift_holders=gift_holders,
                        ppt_cpp_global=ppt_cpp_global,
                        debug=show_debug,
                        debug_overlay=show_layout_overlay,
                    )
                except Exception as e:
                    if show_debug:
                        st.exception(e)
                    else:
                        st.error("Error generating Full Pallet PDF. Enable debug for details.")
                    return
            st.success("Done.")
            st.download_button(
                "⬇ Download Planogram PDF",
                pdf,
                file_name=out_name if out_name.endswith(".pdf") else f"{out_name}.pdf",
                mime="application/pdf",
                use_container_width=True,
            )


if __name__ == "__main__":
    main()
