from __future__ import annotations

import io
import re
from typing import Dict, List, Optional, Tuple

from PIL import Image

from app.shared.models import PptCard, PptSideCards
from app.shared.text_utils import _coerce_int


def load_ppt_cards(pptx_bytes: bytes) -> Dict[str, PptSideCards]:
    """Parse PPTX per-side slides and return PptSideCards with images.

    Structure of each SIDE slide (slides 2-5 = A-D):
      - Top-level GROUP shapes (one per card image, width < 10in) sitting in the top ~1in band
      - TEXT_BOX shapes with "ID #nn" labels positioned just below each group
      - One loose PICTURE shape may appear on slides where one card isn't inside a group

    Each card GROUP contains:
      - The actual card image (variable blob size)
      - The Walmart watermark/logo (constant 71369-byte blob) — skip this

    Algorithm:
      1. Collect all image containers (groups + loose pictures), extract card blob
      2. Collect all TEXT_BOX labels with ID#
      3. Match each label to the nearest image container above it (by cx proximity)
      4. Split into top8 (labels at y < 40% slide height) and side6 (y >= 40%)
      5. Sort top8 left→right by cx; sort side6 by (row-band, cx)
    """
    from pptx import Presentation  # type: ignore
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore

    WATERMARK_BLOB_SIZE = 71369  # Walmart logo — same bytes across all groups

    prs = Presentation(io.BytesIO(pptx_bytes))
    slide_h = float(prs.slide_height)
    slide_w = float(prs.slide_width)

    id_re = re.compile(r"\bID\s*#?\s*[:\-]?\s*(\d{1,8})\b", re.IGNORECASE)
    side_re = re.compile(r"\bSIDE\s*([A-D])\b", re.IGNORECASE)

    def _extract_card_pictures(sh) -> List[dict]:
        """Return the largest non-watermark image blob from a shape or group."""
        pics: List[dict] = []

        def _collect(shapes) -> None:
            for s in shapes:
                st = getattr(s, "shape_type", None)
                if st == MSO_SHAPE_TYPE.PICTURE:
                    b = bytes(s.image.blob)
                    if len(b) == WATERMARK_BLOB_SIZE:
                        continue
                    l = float(getattr(s, "left", 0) or 0)
                    t = float(getattr(s, "top", 0) or 0)
                    w = float(getattr(s, "width", 0) or 0)
                    h = float(getattr(s, "height", 0) or 0)
                    pics.append({
                        "blob": b,
                        "cx": l + w / 2,
                        "cy": t + h / 2,
                        "area": w * h,
                        "byte_size": len(b)
                    })

                elif st == MSO_SHAPE_TYPE.GROUP:
                    _collect(s.shapes)

        st = getattr(sh, "shape_type", None)
        if st == MSO_SHAPE_TYPE.PICTURE:
            b = bytes(sh.image.blob)
            if len(b) != WATERMARK_BLOB_SIZE:
                l = float(getattr(sh, "left", 0) or 0)
                t = float(getattr(sh, "top", 0) or 0)
                w = float(getattr(sh, "width", 0) or 0)
                h = float(getattr(sh, "height", 0) or 0)
                pics.append({
                    "blob": b,
                    "cx": l + w / 2,
                    "cy": t + h / 2,
                    "area": w * h,
                    "byte_size": len(b),
                })
        elif st == MSO_SHAPE_TYPE.GROUP:
            _collect(sh.shapes)

        return pics

    best_by_side: Dict[str, Tuple[int, List[PptCard], List[PptCard]]] = {}

    for slide in prs.slides:
        # ── Detect side letter ──────────────────────────────────────────────
        side_letter: Optional[str] = None
        for sh in slide.shapes:
            m = side_re.search(str(getattr(sh, "text", "") or ""))
            if m:
                side_letter = m.group(1).upper()
                break
        if side_letter is None:
            continue  # skip title / intro / half-pallet slides

        # ── Collect image containers ────────────────────────────────────────
        img_containers: List[dict] = []
        for sh in slide.shapes:
            stype = getattr(sh, "shape_type", None)
            l = float(getattr(sh, "left", 0) or 0)
            t = float(getattr(sh, "top", 0) or 0)
            w = float(getattr(sh, "width", 0) or 0)
            h = float(getattr(sh, "height", 0) or 0)

            if stype == MSO_SHAPE_TYPE.GROUP:
                if w > slide_w * 0.80:
                    continue  # skip full-slide background frame
                candidates = _extract_card_pictures(sh)
                if not candidates:
                    continue
            elif stype == MSO_SHAPE_TYPE.PICTURE:
                candidates = _extract_card_pictures(sh)
                if not candidates:
                    continue
            else:
                continue

            img_containers.append({
                "cx": l + w / 2,
                "cy": t + h / 2,
                "bottom": t + h,
                "candidates": candidates,
            })

        # ── Collect labels ──────────────────────────────────────────────────
        labels: List[dict] = []
        for sh in slide.shapes:
            txt = str(getattr(sh, "text", "") or "").strip()
            m = id_re.search(txt)
            if not m:
                continue
            card_id = m.group(1)
            lines = [ln.strip() for ln in txt.splitlines() if ln.strip()]
            id_idx = next((i for i, ln in enumerate(lines) if id_re.search(ln)), 0)
            title = " ".join(lines[:id_idx]).strip()
            l = float(getattr(sh, "left", 0) or 0)
            t = float(getattr(sh, "top", 0) or 0)
            w = float(getattr(sh, "width", 0) or 0)
            h = float(getattr(sh, "height", 0) or 0)
            labels.append({
                "card_id": card_id,
                "title": title,
                "cx": l + w / 2,
                "top": t,
                "cy": t + h / 2,
            })

        if not labels:
            continue

        # ── Match each label to the nearest image above it ──────────────────
        # Process in reading order (top→bottom, left→right) so greedy works well.
        used: set = set()
        label_to_img: Dict[str, Optional[dict]] = {}

        for lab in sorted(labels, key=lambda d: (round(d["top"] / (slide_h * 0.05)), d["cx"])):
            best_i: Optional[int] = None
            best_score = float("inf")
            for i, img in enumerate(img_containers):
                if i in used:
                    continue
                # Image must be above (or at most 0.5in below) the label top
                if img["bottom"] > lab["top"] + slide_h * 0.07:
                    continue
                dx = abs(lab["cx"] - img["cx"])
                dy = max(0.0, lab["top"] - img["bottom"])
                score = dx + 0.3 * dy
                if score < best_score:
                    best_score = score
                    best_i = i
            if best_i is not None:
                used.add(best_i)
                label_to_img[lab["card_id"]] = img_containers[best_i]
            else:
                label_to_img[lab["card_id"]] = None

        # ── Split into top8 and side6 by y-position ─────────────────────────
        # Labels for the main top row appear at ~30-36% of slide height.
        # Labels for the side panel appear at ~57-93% of slide height.
        TOP_THRESH = slide_h * 0.42
        top_labels = sorted(
            [lb for lb in labels if lb["top"] < TOP_THRESH], key=lambda d: d["cx"]
        )
        side_labels = sorted(
            [lb for lb in labels if lb["top"] >= TOP_THRESH],
            key=lambda d: (round(d["top"] / (slide_h * 0.12)), d["cx"]),
        )

        def _make_card(lab: dict) -> PptCard:
            img_entry = label_to_img.get(lab["card_id"])
            img_bytes: Optional[bytes] = None
            img_ext: Optional[str] = None
            if img_entry:
                candidates = img_entry.get("candidates", [])
                chosen: Optional[dict] = None

                if candidates:
                    chosen = min(
                        candidates,
                        key=lambda c: (abs(c["cx"] - lab["cx"]), abs(c["cy"] - lab["cy"]), -c["area"], -c["byte_size"]),
                    )
                # Normalise to PNG
                if chosen:
                    raw = chosen["blob"]
                    try:
                        from io import BytesIO as _BIO
                        im = Image.open(_BIO(raw)).convert("RGBA")
                        out = _BIO()
                        im.save(out, format="PNG")
                        img_bytes = out.getvalue()
                        img_ext = "png"
                    except Exception:
                        img_bytes = raw
                        img_ext = "png"

            return PptCard(
                card_id=lab["card_id"],
                title=lab["title"],
                image_bytes=img_bytes,
                image_ext=img_ext,
            )

        top8 = [_make_card(lb) for lb in top_labels[:8]]
        side6 = [_make_card(lb) for lb in side_labels[:6]]

        total = len(top8) + len(side6)
        prev = best_by_side.get(side_letter)
        if prev is None or total > prev[0]:
            best_by_side[side_letter] = (total, top8, side6)

    parsed: Dict[str, PptSideCards] = {}
    for side in "ABCD":
        entry = best_by_side.get(side)
        if entry:
            _, top8, side6 = entry
        else:
            top8, side6 = [], []
        parsed[side] = PptSideCards(side=side, top8=top8, side6=side6)

    return parsed



def validate_ppt_side_cards(ppt_cards: Dict[str, PptSideCards]) -> List[str]:
    issues: List[str] = []
    for side in "ABCD":
        side_cards = ppt_cards.get(side, PptSideCards(side=side, top8=[], side6=[]))
        if len(side_cards.top8) != 8 or len(side_cards.side6) != 6:
            issues.append(
                f"SIDE {side}: found top8={len(side_cards.top8)} and side6={len(side_cards.side6)} (expected 8 and 6)."
            )
    return issues
